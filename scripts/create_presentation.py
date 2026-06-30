#!/usr/bin/env python3
"""
Gera a apresentação PowerPoint do projeto:
  Classificação de Radiografias Torácicas com ResNet-50 e Grad-CAM++

Pré-requisitos:
  pip install python-pptx
  (figuras em figures/ e métricas em results/ já geradas)

Uso:
  python scripts/create_presentation.py
  python scripts/create_presentation.py --out apresentacao.pptx
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns
from PIL import Image
from sklearn.metrics import auc, roc_curve

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
RESULTS_DIR = ROOT / "results"
FIGURES_DIR = ROOT / "figures"


# ──────────────────────────────────────────────────────────────────────────────
# Paleta de cores
# ──────────────────────────────────────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x10, 0x3A, 0x6E)   # azul escuro (títulos / fundo)
C_MED_BLUE   = RGBColor(0x1A, 0x6C, 0xB0)   # azul médio
C_LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)   # azul claro (caixas)
C_ACCENT     = RGBColor(0xE8, 0x6A, 0x2B)   # laranja (destaques)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
C_MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
C_GREEN      = RGBColor(0x2C, 0xA0, 0x2C)
C_RED        = RGBColor(0xD6, 0x27, 0x28)
C_YELLOW     = RGBColor(0xFF, 0xBF, 0x00)


# ──────────────────────────────────────────────────────────────────────────────
# Helpers de layout
# ──────────────────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _rgb(r: RGBColor):
    from pptx.dml.color import RGBColor as RC
    return RC(r[0], r[1], r[2])


def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_width_pt: float = 0.0):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text: str, left, top, width, height,
                 font_size: int = 18, bold: bool = False, italic: bool = False,
                 color: RGBColor = C_DARK_GRAY, align=PP_ALIGN.LEFT,
                 word_wrap: bool = True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_list(slide, items: list[str], left, top, width, height,
                    font_size: int = 16, color: RGBColor = C_DARK_GRAY,
                    bullet_char: str = "•  ", bold_first: bool = False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first and idx == 0:
            run.font.bold = True
    return txb


def fig_to_bytesio(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def image_to_bytesio(img_path: str | Path) -> io.BytesIO:
    img = Image.open(img_path).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ──────────────────────────────────────────────────────────────────────────────
# Geração de figuras inline
# ──────────────────────────────────────────────────────────────────────────────

def make_confusion_matrix_png(cm: np.ndarray, title: str) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(5, 4))
    sns.heatmap(cm, annot=True, fmt="d", cmap="Blues", cbar=False, ax=ax,
                annot_kws={"size": 16, "weight": "bold"})
    ax.set_title(title, fontsize=13, fontweight="bold", pad=10)
    ax.set_xlabel("Predito", fontsize=11)
    ax.set_ylabel("Verdadeiro", fontsize=11)
    ax.set_xticklabels(["NORMAL", "PNEUMONIA"], fontsize=10)
    ax.set_yticklabels(["NORMAL", "PNEUMONIA"], fontsize=10, rotation=0)
    fig.tight_layout()
    return fig_to_bytesio(fig)


def make_roc_png(predictions: list[dict], title: str) -> io.BytesIO:
    y_true = np.array([int(p["true_label"]) for p in predictions])
    y_score = np.array([float(p["score_pneumonia"]) for p in predictions])
    fpr, tpr, _ = roc_curve(y_true, y_score)
    auc_val = auc(fpr, tpr)
    fig, ax = plt.subplots(figsize=(5, 4))
    ax.plot([0, 1], [0, 1], "--", color="gray", label="Aleatório")
    ax.plot(fpr, tpr, color="navy", lw=2, label=f"AUC = {auc_val:.4f}")
    ax.set_title(title, fontsize=13, fontweight="bold")
    ax.set_xlabel("False Positive Rate", fontsize=11)
    ax.set_ylabel("True Positive Rate", fontsize=11)
    ax.legend(loc="lower right", fontsize=11)
    ax.fill_between(fpr, tpr, alpha=0.1, color="navy")
    fig.tight_layout()
    return fig_to_bytesio(fig)


def make_error_bar_png(counts: dict, title: str) -> io.BytesIO:
    labels = ["focused", "partial", "distracted"]
    values = [int(counts.get(k, 0)) for k in labels]
    colors = ["#2ca02c", "#ffbf00", "#d62728"]
    fig, ax = plt.subplots(figsize=(5, 4))
    bars = ax.bar(labels, values, color=colors, edgecolor="white", linewidth=1.2)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                str(val), ha="center", va="bottom", fontsize=13, fontweight="bold")
    ax.set_title(title, fontsize=13, fontweight="bold")
    ax.set_xlabel("Categoria de atenção", fontsize=11)
    ax.set_ylabel("Contagem", fontsize=11)
    ax.set_ylim(0, max(values) * 1.25 if max(values) > 0 else 5)
    fig.tight_layout()
    return fig_to_bytesio(fig)


def make_metrics_comparison_png(m_head: dict, m_full: dict) -> io.BytesIO:
    metrics_names = ["Acurácia (%)", "AUC-ROC", "F1", "Precisão", "Recall"]
    head_vals = [
        m_head["test"]["accuracy_percent"],
        round(m_head["test"]["auc_roc"] * 100, 2),
        round(m_head["test"]["f1"] * 100, 2),
        round(m_head["test"]["precision"] * 100, 2),
        round(m_head["test"]["recall"] * 100, 2),
    ]
    full_vals = [
        m_full["test"]["accuracy_percent"],
        round(m_full["test"]["auc_roc"] * 100, 2),
        round(m_full["test"]["f1"] * 100, 2),
        round(m_full["test"]["precision"] * 100, 2),
        round(m_full["test"]["recall"] * 100, 2),
    ]
    x = np.arange(len(metrics_names))
    width = 0.35
    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars1 = ax.bar(x - width / 2, head_vals, width, label="Head Only", color="#5B9BD5", edgecolor="white")
    bars2 = ax.bar(x + width / 2, full_vals, width, label="Full Finetune", color="#ED7D31", edgecolor="white")
    for bar in list(bars1) + list(bars2):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                f"{bar.get_height():.1f}", ha="center", va="bottom", fontsize=8)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics_names, fontsize=10)
    ax.set_ylabel("Valor (%)", fontsize=11)
    ax.set_title("Comparação de Métricas no Conjunto de Teste", fontsize=12, fontweight="bold")
    ax.legend(fontsize=10)
    ax.set_ylim(0, 115)
    ax.yaxis.grid(True, linestyle="--", alpha=0.5)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytesio(fig)


def make_timeline_png() -> io.BytesIO:
    weeks = [
        ("S1\n10–13/abr", "Fundação\nteórica", True),
        ("S2\n14–20/abr", "Ambiente\ne dados", True),
        ("S3\n21/abr–1/mai", "Modelo\nbase", True),
        ("S4\n2–8/mai", "Análise\nde erros", True),
        ("S5\n9–15/mai", "Preparação\napresentação", True),
        ("S6\n16–22/mai", "XAI-guided\nfine-tuning", False),
        ("S7\n23–29/mai", "Avaliação\nfinal", False),
    ]
    fig, ax = plt.subplots(figsize=(11, 2.8))
    ax.axis("off")
    n = len(weeks)
    for i, (label, task, done) in enumerate(weeks):
        x = i / (n - 1)
        color = "#2ca02c" if done else "#ED7D31"
        ax.plot(x, 0.5, "o", markersize=18, color=color, zorder=3)
        if done:
            ax.text(x, 0.5, "✓", ha="center", va="center", fontsize=10,
                    color="white", fontweight="bold", zorder=4)
        if i < n - 1:
            next_x = (i + 1) / (n - 1)
            ax.plot([x, next_x], [0.5, 0.5], "-", color="#cccccc", lw=2, zorder=1)
        ax.text(x, 0.12, label, ha="center", va="top", fontsize=8,
                color="#555555", multialignment="center")
        ax.text(x, 0.85, task, ha="center", va="bottom", fontsize=8,
                color="#333333", fontweight="bold", multialignment="center")
    from matplotlib.patches import Patch
    legend_elements = [
        Patch(facecolor="#2ca02c", label="Concluída"),
        Patch(facecolor="#ED7D31", label="Em andamento / pendente"),
    ]
    ax.legend(handles=legend_elements, loc="lower right", fontsize=9,
              bbox_to_anchor=(1.0, -0.1))
    ax.set_xlim(-0.08, 1.08)
    ax.set_ylim(-0.05, 1.3)
    ax.set_title("Cronograma do Projeto", fontsize=12, fontweight="bold", pad=5)
    fig.tight_layout()
    return fig_to_bytesio(fig)


def make_architecture_png() -> io.BytesIO:
    """Diagrama de blocos da ResNet-50 usando coordenadas de dados (sem transAxes)."""
    from matplotlib.patches import FancyArrowPatch, Patch
    from matplotlib.patches import FancyBboxPatch

    fig, ax = plt.subplots(figsize=(12, 3.6))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 3)
    ax.axis("off")

    # (label, x_center, cor_fundo)
    boxes = [
        ("Input\n224×224×3",       1.0,  "#AED6F1"),
        ("Conv1\n+ Pool\n(frozen)", 3.0,  "#D5E8D4"),
        ("Layer1–2\n(frozen)",      5.0,  "#D5E8D4"),
        ("Layer3\n(full_ft)",       7.0,  "#FFE6CC"),
        ("Layer4\n(full_ft)",       9.0,  "#FFE6CC"),
        ("AvgPool\n+ Flatten",     11.0,  "#F8CECC"),
        ("FC\n2 classes",          13.0,  "#F8CECC"),
    ]

    BOX_W, BOX_H = 1.6, 1.6
    BOX_Y = 0.7   # bottom y

    for label, cx, color in boxes:
        rect = FancyBboxPatch(
            (cx - BOX_W / 2, BOX_Y), BOX_W, BOX_H,
            boxstyle="round,pad=0.05",
            facecolor=color, edgecolor="#555555", linewidth=1.4,
        )
        ax.add_patch(rect)
        ax.text(cx, BOX_Y + BOX_H / 2, label,
                ha="center", va="center", fontsize=8.5, fontweight="bold",
                multialignment="center", color="#222222")

    # Setas entre boxes (em coordenadas de dados)
    for i in range(len(boxes) - 1):
        x_start = boxes[i][1]   + BOX_W / 2 + 0.08
        x_end   = boxes[i+1][1] - BOX_W / 2 - 0.08
        y_mid   = BOX_Y + BOX_H / 2
        arrow = FancyArrowPatch(
            (x_start, y_mid), (x_end, y_mid),
            arrowstyle="-|>", color="#444444",
            mutation_scale=12, lw=1.5,
        )
        ax.add_patch(arrow)

    # Legenda
    legend_handles = [
        Patch(facecolor="#D5E8D4", edgecolor="#555", label="Congelado (ambas as estratégias)"),
        Patch(facecolor="#FFE6CC", edgecolor="#555", label="Descongelado em full_finetune"),
        Patch(facecolor="#F8CECC", edgecolor="#555", label="Treinado em ambas"),
    ]
    ax.legend(handles=legend_handles, loc="lower center", fontsize=8.5,
              bbox_to_anchor=(0.5, -0.08), ncol=3, framealpha=0.9)

    ax.set_title("Arquitetura ResNet-50 — Estratégias de Transfer Learning",
                 fontsize=11, fontweight="bold", pad=6)
    fig.tight_layout()
    return fig_to_bytesio(fig)


def make_xai_concept_png() -> io.BytesIO:
    fig, axes = plt.subplots(1, 3, figsize=(10, 3.5))
    for ax in axes:
        ax.axis("off")
    # Simula imagem X-ray
    np.random.seed(42)
    base = np.zeros((100, 100, 3), dtype=np.float32)
    for _ in range(2000):
        cx, cy = np.random.randint(20, 80), np.random.randint(25, 75)
        r = np.random.randint(2, 8)
        v = np.random.uniform(0.2, 0.8)
        yy, xx = np.ogrid[:100, :100]
        mask = (xx - cx)**2 + (yy - cy)**2 < r**2
        base[mask] = np.clip(base[mask] + v, 0, 1)
    gray = (base[:, :, 0] * 0.5 + 0.2).clip(0, 1)
    axes[0].imshow(gray, cmap="gray", vmin=0, vmax=1)
    axes[0].set_title("Radiografia Original", fontsize=10, fontweight="bold")
    axes[0].axis("off")
    # CAM "distracted" (attention outside lung)
    cam_bad = np.zeros((100, 100))
    cam_bad[:20, :] = np.random.rand(20, 100) * 0.8  # top corner
    cam_bad[:, :15] = np.random.rand(100, 15) * 0.8  # left border
    cam_bad = cam_bad / (cam_bad.max() + 1e-8)
    overlay_bad = gray[:, :, np.newaxis] * np.array([0.7, 0.7, 0.7])
    jet = plt.cm.jet(cam_bad)[:, :, :3]
    overlay_bad = np.clip(overlay_bad + 0.5 * jet, 0, 1)
    axes[1].imshow(overlay_bad)
    axes[1].add_patch(plt.Rectangle((15, 10), 70, 80, fill=False,
                                    edgecolor="lime", linewidth=2, linestyle="--"))
    axes[1].set_title("CAM 'Distraída'\n(atenção fora do pulmão)", fontsize=9, fontweight="bold",
                      color="#D62728")
    axes[1].axis("off")
    # CAM "focused" (attention inside lung)
    yy, xx = np.mgrid[:100, :100]
    cam_good = np.exp(-((xx - 50)**2 / 600 + (yy - 50)**2 / 800))
    cam_good = (cam_good - cam_good.min()) / (cam_good.max() - cam_good.min())
    overlay_good = gray[:, :, np.newaxis] * np.array([0.7, 0.7, 0.7])
    jet2 = plt.cm.jet(cam_good)[:, :, :3]
    overlay_good = np.clip(overlay_good + 0.5 * jet2, 0, 1)
    axes[2].imshow(overlay_good)
    axes[2].add_patch(plt.Rectangle((15, 10), 70, 80, fill=False,
                                    edgecolor="lime", linewidth=2, linestyle="--"))
    axes[2].set_title("CAM 'Focada'\n(atenção na região pulmonar)", fontsize=9, fontweight="bold",
                      color="#2CA02C")
    axes[2].axis("off")
    plt.suptitle("Grad-CAM++: Qualidade da Atenção do Modelo",
                 fontsize=11, fontweight="bold", y=1.02)
    fig.tight_layout()
    return fig_to_bytesio(fig)


# ──────────────────────────────────────────────────────────────────────────────
# Montagem dos slides
# ──────────────────────────────────────────────────────────────────────────────

def make_header_bar(slide, title_text: str, subtitle_text: str = ""):
    """Barra superior azul escura com título e subtítulo."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=C_DARK_BLUE)
    add_text_box(slide, title_text,
                 Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.7),
                 font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text_box(slide, subtitle_text,
                     Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.4),
                     font_size=14, bold=False, color=C_LIGHT_BLUE, align=PP_ALIGN.LEFT)
    # Linha colorida embaixo do header
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.04), fill=C_ACCENT)


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Fundo dividido
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK_BLUE)
    add_rect(slide, Inches(7.5), 0, Inches(5.83), SLIDE_H, fill=C_MED_BLUE)
    add_rect(slide, 0, Inches(6.5), SLIDE_W, Inches(1.0), fill=C_ACCENT)

    # Título
    add_text_box(slide, "Classificação de Radiografias\nTorácicas com ResNet-50\ne Grad-CAM++",
                 Inches(0.5), Inches(1.2), Inches(6.8), Inches(3.0),
                 font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # Subtítulo
    add_text_box(slide, "XAI-Guided Fine-Tuning para Diagnóstico de Pneumonia",
                 Inches(0.5), Inches(4.2), Inches(6.8), Inches(0.8),
                 font_size=16, bold=False, color=C_LIGHT_BLUE, align=PP_ALIGN.LEFT)
    # Info
    add_text_box(slide, "Disciplina: Visão Computacional\nUFF — 2026",
                 Inches(0.5), Inches(5.0), Inches(6.8), Inches(0.8),
                 font_size=13, bold=False, color=C_MID_GRAY, align=PP_ALIGN.LEFT)
    # Divisor
    add_rect(slide, Inches(0.5), Inches(5.9), Inches(4.0), Inches(0.04), fill=C_ACCENT)


def slide_motivation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Motivação e Contexto",
                    "Por que classificar radiografias torácicas automaticamente?")
    # Coluna esquerda
    add_rect(slide, Inches(0.3), Inches(1.45), Inches(5.9), Inches(5.6),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Desafio Clínico",
                 Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                 font_size=15, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "Pneumonia causa ~2,56 milhões de mortes/ano (OMS)",
        "Radiografia torácica é o exame diagnóstico mais acessível",
        "Interpretação manual é subjetiva e demorada",
        "Faltam radiologistas em regiões remotas e países em desenvolvimento",
    ], Inches(0.4), Inches(2.0), Inches(5.8), Inches(3.5), font_size=14)

    # Coluna direita
    add_rect(slide, Inches(6.5), Inches(1.45), Inches(6.5), Inches(5.6),
             fill=RGBColor(0xF0, 0xF0, 0xF0), line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Proposta Computacional",
                 Inches(6.7), Inches(1.5), Inches(6.0), Inches(0.5),
                 font_size=15, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "Transfer learning com ResNet-50 pré-treinada no ImageNet",
        "Classificação binária: NORMAL vs PNEUMONIA",
        "Explicabilidade com Grad-CAM++ para auditoria do modelo",
        "Fine-tuning guiado por mapas de atenção (XAI-guided)",
        "Pipeline reproduzível e avaliado com métricas clínicas",
    ], Inches(6.6), Inches(2.0), Inches(6.2), Inches(4.0), font_size=14)


def slide_dataset(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Dataset: Chest X-Ray Images (Pneumonia)",
                    "Kaggle — Paul Mooney · 5,856 radiografias")

    # Stats boxes
    stats = [
        ("5,856\ntotal", C_DARK_BLUE, C_WHITE),
        ("4,434\ntreino", C_MED_BLUE, C_WHITE),
        ("624\nteste", C_ACCENT, C_WHITE),
        ("2 classes\nNORMAL / PNEUMONIA", RGBColor(0x1E, 0x8B, 0x4C), C_WHITE),
    ]
    bw = Inches(2.8)
    for i, (txt, bg, fg) in enumerate(stats):
        x = Inches(0.3 + i * 3.17)
        add_rect(slide, x, Inches(1.45), bw, Inches(1.1), fill=bg)
        add_text_box(slide, txt, x, Inches(1.45), bw, Inches(1.1),
                     font_size=16, bold=True, color=fg, align=PP_ALIGN.CENTER)

    # Desbalanceamento
    add_rect(slide, Inches(0.3), Inches(2.7), Inches(5.9), Inches(4.1),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Distribuição do Treino",
                 Inches(0.5), Inches(2.75), Inches(5.5), Inches(0.45),
                 font_size=14, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "NORMAL:    1,341 imagens  (30%)",
        "PNEUMONIA: 3,305 imagens  (74%)",
        "",
        "Desbalanceamento  ~1:2,9",
        "→ WeightedRandomSampler para equilíbrio no treino",
        "→ CrossEntropyLoss com pesos de classe inversamente\n   proporcionais à frequência",
    ], Inches(0.5), Inches(3.2), Inches(5.5), Inches(3.3), font_size=13)

    # Pré-processamento
    add_rect(slide, Inches(6.5), Inches(2.7), Inches(6.5), Inches(4.1),
             fill=RGBColor(0xF0, 0xF0, 0xF0), line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Pré-processamento",
                 Inches(6.7), Inches(2.75), Inches(6.0), Inches(0.45),
                 font_size=14, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "Redimensionamento: 224×224 px",
        "Normalização: média e σ do ImageNet",
        "Augmentation no treino:",
        "    • Flip horizontal aleatório",
        "    • Rotação aleatória (±10°)",
        "Validação: split de 15% do treino",
    ], Inches(6.6), Inches(3.2), Inches(6.2), Inches(3.3), font_size=13)


def slide_architecture(prs: Presentation, arch_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Arquitetura e Estratégias de Transfer Learning",
                    "ResNet-50 pré-treinada no ImageNet → classificação binária")

    slide.shapes.add_picture(arch_buf, Inches(0.3), Inches(1.35), Inches(12.7), Inches(3.1))

    # Tabela comparativa
    add_rect(slide, Inches(0.3), Inches(4.6), Inches(12.7), Inches(2.55),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Comparação das Estratégias",
                 Inches(0.5), Inches(4.65), Inches(12.0), Inches(0.4),
                 font_size=14, bold=True, color=C_DARK_BLUE)

    headers = ["Estratégia", "Camadas treináveis", "LR backbone", "LR FC", "Épocas"]
    h_text = "    ".join(f"{h:20s}" for h in headers)
    add_text_box(slide, h_text,
                 Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.35),
                 font_size=11, bold=True, color=C_DARK_BLUE)
    rows = [
        ["head_only", "FC apenas", "— (congelado)", "1e-3", "5"],
        ["full_finetune", "layer3 + layer4 + FC", "1e-5 (÷10)", "1e-4", "10"],
    ]
    for i, row in enumerate(rows):
        r_text = "    ".join(f"{v:20s}" for v in row)
        bg = C_WHITE if i % 2 == 0 else RGBColor(0xE8, 0xF4, 0xFF)
        add_text_box(slide, r_text,
                     Inches(0.5), Inches(5.48 + i * 0.35), Inches(12.3), Inches(0.32),
                     font_size=11, bold=False, color=C_DARK_GRAY)


def slide_strategies(prs: Presentation, m_head: dict, m_full: dict) -> None:
    """Slide comparando Head Only vs Full Finetune — estratégias de transfer learning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Head Only vs Full Finetune — Estratégias de Transfer Learning",
                    "Como e por que treinamos de formas diferentes")

    # ── Painel HEAD ONLY (esquerda) ──────────────────────────────────────────
    add_rect(slide, Inches(0.25), Inches(1.35), Inches(6.1), Inches(4.45),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=1.2)
    add_text_box(slide, "HEAD ONLY",
                 Inches(0.4), Inches(1.4), Inches(4.0), Inches(0.42),
                 font_size=16, bold=True, color=C_DARK_BLUE)
    add_text_box(slide, "Treinar só a cabeça da rede",
                 Inches(0.4), Inches(1.82), Inches(5.7), Inches(0.32),
                 font_size=12, bold=False, color=C_MED_BLUE)

    # Mini diagrama de camadas
    frozen_color = RGBColor(0x2C, 0xA0, 0x2C)   # verde = congelado
    active_color = RGBColor(0xE8, 0x6A, 0x2B)   # laranja = ativo
    layers_head = [
        ("Conv1–Layer2", frozen_color, "congelado"),
        ("Layer3",       frozen_color, "congelado"),
        ("Layer4",       frozen_color, "congelado"),
        ("AvgPool",      frozen_color, "congelado"),
        ("FC  →  2",     active_color, "treina"),
    ]
    for k, (lbl, col, tag) in enumerate(layers_head):
        bx = Inches(0.4)
        by = Inches(2.2 + k * 0.44)
        add_rect(slide, bx, by, Inches(3.0), Inches(0.36), fill=col)
        add_text_box(slide, lbl, bx, by, Inches(3.0), Inches(0.36),
                     font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, tag,
                     Inches(3.5), by, Inches(2.5), Inches(0.36),
                     font_size=9, bold=False, color=C_MID_GRAY)

    add_bullet_list(slide, [
        f"Parâmetros ajustados: 4.098  (2048×2 + 2)",
        f"LR FC: 1e-3  |  Backbone: congelado",
        f"Épocas: 5  (converge rápido)",
        f"Sem risco de catástrofe do esquecimento",
        f"AUC teste: {m_head['test']['auc_roc']:.4f}  |  Recall: {m_head['test']['recall']:.4f}",
    ], Inches(0.4), Inches(4.45), Inches(5.7), Inches(1.2), font_size=11)

    # ── Painel FULL FINETUNE (direita) ───────────────────────────────────────
    add_rect(slide, Inches(6.65), Inches(1.35), Inches(6.45), Inches(4.45),
             fill=RGBColor(0xFF, 0xF3, 0xE8), line=C_ACCENT, line_width_pt=1.2)
    add_text_box(slide, "FULL FINETUNE",
                 Inches(6.8), Inches(1.4), Inches(4.5), Inches(0.42),
                 font_size=16, bold=True, color=C_ACCENT)
    add_text_box(slide, "Adaptar o backbone ao domínio médico",
                 Inches(6.8), Inches(1.82), Inches(6.0), Inches(0.32),
                 font_size=12, bold=False, color=C_ACCENT)

    layers_full = [
        ("Conv1–Layer2", frozen_color, "congelado"),
        ("Layer3",       active_color, "LR 1e-5"),
        ("Layer4",       active_color, "LR 1e-5"),
        ("AvgPool",      frozen_color, "congelado"),
        ("FC  →  2",     active_color, "LR 1e-4"),
    ]
    for k, (lbl, col, tag) in enumerate(layers_full):
        bx = Inches(6.8)
        by = Inches(2.2 + k * 0.44)
        add_rect(slide, bx, by, Inches(3.0), Inches(0.36), fill=col)
        add_text_box(slide, lbl, bx, by, Inches(3.0), Inches(0.36),
                     font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, tag,
                     Inches(9.9), by, Inches(3.0), Inches(0.36),
                     font_size=9, bold=False, color=C_MID_GRAY)

    add_bullet_list(slide, [
        f"Parâmetros ajustados: ~22 M  (Layer3=7M + Layer4=15M + FC)",
        f"LR FC: 1e-4  |  Backbone: 1e-5  (10× menor)",
        f"Épocas: 10  (blocos profundos precisam adaptar)",
        f"LR baixo no backbone evita catástrofe do esquecimento",
        f"AUC teste: {m_full['test']['auc_roc']:.4f}  |  Recall: {m_full['test']['recall']:.4f}",
    ], Inches(6.8), Inches(4.45), Inches(6.0), Inches(1.2), font_size=11)

    # ── Barra de comparação inferior ─────────────────────────────────────────
    add_rect(slide, Inches(0.25), Inches(5.9), Inches(12.85), Inches(1.5),
             fill=C_DARK_BLUE)
    add_text_box(slide, "Comparação direta no conjunto de teste",
                 Inches(0.4), Inches(5.92), Inches(12.5), Inches(0.35),
                 font_size=12, bold=True, color=C_WHITE)

    comparisons = [
        ("Parâmetros\ntreináveis", "4.098", "~22 M"),
        ("Épocas", "5", "10"),
        ("AUC-ROC", f"{m_head['test']['auc_roc']:.4f}", f"{m_full['test']['auc_roc']:.4f}  ★"),
        ("Recall", f"{m_head['test']['recall']:.4f}", f"{m_full['test']['recall']:.4f}  ★"),
        ("FP / FN", "28 / 52", "80 / 4"),
    ]
    col_w = 12.85 / len(comparisons)
    for j, (label, v_head, v_full) in enumerate(comparisons):
        cx = Inches(0.25 + j * col_w)
        cw = Inches(col_w - 0.05)
        add_text_box(slide, label,
                     cx, Inches(6.28), cw, Inches(0.3),
                     font_size=9, bold=True, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
        add_rect(slide, cx, Inches(6.6), cw, Inches(0.28), fill=C_MED_BLUE)
        add_text_box(slide, v_head, cx, Inches(6.6), cw, Inches(0.28),
                     font_size=10, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, cx, Inches(6.9), cw, Inches(0.28),
                 fill=RGBColor(0xFF, 0xC0, 0x80))
        add_text_box(slide, v_full, cx, Inches(6.9), cw, Inches(0.28),
                     font_size=10, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_results_table(prs: Presentation, m_head: dict, m_full: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Resultados — Comparação de Métricas",
                    "Conjunto de teste: 624 imagens (234 NORMAL / 390 PNEUMONIA)")

    # Header da tabela
    cols = ["Modelo", "Acurácia", "Erro", "AUC-ROC", "F1", "Precisão", "Recall"]
    widths = [2.6, 1.5, 1.3, 1.5, 1.3, 1.5, 1.5]
    xs = [0.3]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)

    header_y = 1.45
    for j, (col, w, x) in enumerate(zip(cols, widths, xs)):
        add_rect(slide, Inches(x), Inches(header_y), Inches(w), Inches(0.5), fill=C_DARK_BLUE)
        add_text_box(slide, col, Inches(x), Inches(header_y), Inches(w), Inches(0.5),
                     font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    def fmt(v, is_pct=True):
        if is_pct:
            return f"{v:.2f}%"
        return f"{v:.4f}"

    rows_data = [
        ("Head Only", m_head["test"], C_LIGHT_BLUE),
        ("Full Finetune", m_full["test"], RGBColor(0xFF, 0xF0, 0xE0)),
    ]
    for i, (name, t, bg) in enumerate(rows_data):
        row_y = header_y + 0.5 + i * 0.55
        row_vals = [
            name,
            fmt(t["accuracy_percent"]),
            fmt(t["error_percent"]),
            f"{t['auc_roc']:.4f}",
            f"{t['f1']:.4f}",
            f"{t['precision']:.4f}",
            f"{t['recall']:.4f}",
        ]
        for j, (v, w, x) in enumerate(zip(row_vals, widths, xs)):
            add_rect(slide, Inches(x), Inches(row_y), Inches(w), Inches(0.52), fill=bg)
            color = C_ACCENT if j == 3 and i == 1 else C_DARK_GRAY
            bold = j == 3 and i == 1
            add_text_box(slide, v, Inches(x + 0.05), Inches(row_y), Inches(w - 0.1), Inches(0.52),
                         font_size=13, bold=bold, color=color, align=PP_ALIGN.CENTER)

    # Destaques
    add_rect(slide, Inches(0.3), Inches(3.1), Inches(12.7), Inches(3.9),
             fill=RGBColor(0xF5, 0xF5, 0xF5), line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Destaques",
                 Inches(0.5), Inches(3.15), Inches(12.0), Inches(0.4),
                 font_size=14, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        f"Full finetune: AUC de {m_full['test']['auc_roc']:.4f} (+{m_full['test']['auc_roc']-m_head['test']['auc_roc']:.4f} vs head_only) — melhor discriminação",
        f"Full finetune: Recall de {m_full['test']['recall']:.4f} — crucial para diagnóstico médico (menos FN = menos pneumonias perdidas)",
        f"Head only: Precisão de {m_head['test']['precision']:.4f} — menos falsos positivos (mais conservador)",
        f"Full finetune: 80 FP vs 4 FN — o modelo erra mais para o lado da pneumonia (comportamento seguro clinicamente)",
        "Validação mostrou overfitting parcial: AUC val=0.9963 vs AUC teste=0.9693 no full_finetune",
    ], Inches(0.5), Inches(3.55), Inches(12.4), Inches(3.3), font_size=13)


def slide_confusion_roc(prs: Presentation, cm_buf: io.BytesIO, roc_buf: io.BytesIO,
                        m_full: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Matriz de Confusão e Curva ROC",
                    "Modelo Full Finetune — conjunto de teste (624 imagens)")

    slide.shapes.add_picture(cm_buf, Inches(0.3), Inches(1.4), Inches(5.5), Inches(4.4))
    slide.shapes.add_picture(roc_buf, Inches(6.0), Inches(1.4), Inches(5.5), Inches(4.4))

    # Notas
    add_rect(slide, Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.2),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    cm = m_full["test"]["confusion_matrix"]
    tn, fp, fn, tp = cm[0][0], cm[0][1], cm[1][0], cm[1][1]
    add_text_box(slide,
                 f"  TN={tn}  FP={fp}  FN={fn}  TP={tp}   |   "
                 f"AUC-ROC={m_full['test']['auc_roc']:.4f}   "
                 f"F1={m_full['test']['f1']:.4f}   "
                 f"Precisão={m_full['test']['precision']:.4f}   "
                 f"Recall={m_full['test']['recall']:.4f}",
                 Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.6),
                 font_size=13, bold=False, color=C_DARK_GRAY)
    add_text_box(slide,
                 "  FP: normais classificadas como pneumonia  |  FN: pneumonias classificadas como normal",
                 Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
                 font_size=12, bold=False, color=C_MID_GRAY)


def slide_gradcam_intro(prs: Presentation, compare_img_path: Path | None,
                        xai_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Explicabilidade: Grad-CAM++",
                    "Gradient-weighted Class Activation Mapping — camada layer4[-1]")

    add_bullet_list(slide, [
        "Grad-CAM++ gera mapas de calor mostrando ONDE o modelo 'olhou' para tomar a decisão",
        "Camada alvo: model.layer4[-1]  (última camada convolucional da ResNet-50)",
        "Overlay com colormap JET (α=0.4) aplicado sobre a radiografia original",
        "624 imagens do teste processadas → TP/TN/FP/FN salvos em results/gradcam/",
    ], Inches(0.4), Inches(1.35), Inches(12.6), Inches(1.5), font_size=14)

    slide.shapes.add_picture(xai_buf, Inches(0.3), Inches(2.85), Inches(12.6), Inches(3.0))

    if compare_img_path and compare_img_path.is_file():
        buf = image_to_bytesio(compare_img_path)
        add_text_box(slide, "Comparação real: head_only vs full_finetune (mesma imagem)",
                     Inches(0.3), Inches(5.9), Inches(12.6), Inches(0.35),
                     font_size=12, bold=True, color=C_DARK_BLUE)
        slide.shapes.add_picture(buf, Inches(2.0), Inches(6.2), Inches(9.0), Inches(1.1))


def slide_error_analysis(prs: Presentation, bar_buf: io.BytesIO, error_json: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Análise de Erros com Grad-CAM++",
                    "FP e FN categorizados por overlap_score fora da máscara pulmonar proxy")

    bc = error_json["bucket_counts"]
    ec = error_json.get("error_category_counts", {})

    # Stats boxes
    stats = [
        (f"TP\n{bc['TP']}", C_GREEN),
        (f"TN\n{bc['TN']}", C_MED_BLUE),
        (f"FP\n{bc['FP']}", C_ACCENT),
        (f"FN\n{bc['FN']}", C_RED),
    ]
    for i, (txt, col) in enumerate(stats):
        x = Inches(0.3 + i * 3.0)
        add_rect(slide, x, Inches(1.4), Inches(2.7), Inches(0.95), fill=col)
        add_text_box(slide, txt, x, Inches(1.4), Inches(2.7), Inches(0.95),
                     font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    slide.shapes.add_picture(bar_buf, Inches(0.3), Inches(2.45), Inches(5.5), Inches(4.1))

    # Explicação
    add_rect(slide, Inches(6.0), Inches(2.45), Inches(7.0), Inches(4.1),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Metodologia da Análise",
                 Inches(6.2), Inches(2.5), Inches(6.6), Inches(0.4),
                 font_size=14, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "Máscara pulmonar proxy: crop central de 60%",
        "overlap_score = % ativação FORA da máscara",
        "",
        "focused:     score < 0.20 — atenção no pulmão",
        "partial:     0.20 ≤ score ≤ 0.50",
        "distracted:  score > 0.50 — atenção espúria",
        "",
        f"→ {ec.get('focused',0)} focused  |  {ec.get('partial',0)} partial  |  {ec.get('distracted',0)} distracted",
        f"→ Hard set (overlap > 0.3): {len(error_json.get('hard_set_overlap_gt_0.3', []))} casos",
        "",
        "80 FP (Normal → Pneumonia)",
        "4  FN (Pneumonia → Normal)",
        "Maioria partial: atenção parcialmente fora do pulmão",
    ], Inches(6.1), Inches(2.9), Inches(6.6), Inches(3.5), font_size=12)


def make_error_grid_png(fp_files: list[Path], fn_files: list[Path]) -> io.BytesIO:
    """Gera um grid matplotlib 4×4 com 8 FP + FN (máx disponível) com overlays Grad-CAM++."""
    fp_sel = fp_files[:8]
    fn_sel = fn_files[:8]
    all_items = [(p, "FP") for p in fp_sel] + [(p, "FN") for p in fn_sel]
    # Garante 16 slots (preenche com None se necessário)
    while len(all_items) < 16:
        all_items.append((None, ""))

    fig, axes = plt.subplots(4, 4, figsize=(12, 11))
    fig.patch.set_facecolor("#1a1a1a")
    for i, ax in enumerate(axes.flatten()):
        ax.axis("off")
        item_path, lbl = all_items[i]
        if item_path is None:
            ax.set_facecolor("#1a1a1a")
            continue
        try:
            img = Image.open(item_path).convert("RGB")
            ax.imshow(np.asarray(img))
            color = "#FF6B35" if lbl == "FP" else "#FF3333"
            label_full = "FP — Normal → Pneumonia" if lbl == "FP" else "FN — Pneumonia → Normal"
            ax.set_title(label_full, fontsize=7.5, color=color,
                         fontweight="bold", pad=2)
        except Exception:
            ax.set_facecolor("#333333")
            ax.text(0.5, 0.5, f"{lbl}\n(sem img)", ha="center", va="center",
                    fontsize=8, color="white", transform=ax.transAxes)

    fig.suptitle("Grid de Erros — Overlays Grad-CAM++  (laranja = FP | vermelho = FN)",
                 fontsize=11, fontweight="bold", color="white", y=0.99)
    fig.tight_layout(rect=[0, 0, 1, 0.97])
    return fig_to_bytesio(fig)


def slide_error_grid(prs: Presentation, grid_pdf_path: Path) -> None:
    """Slide com grid real de FP/FN gerado a partir dos PNGs Grad-CAM++."""
    fp_dir = ROOT / "results" / "gradcam" / "FP"
    fn_dir = ROOT / "results" / "gradcam" / "FN"
    fp_files = sorted(fp_dir.glob("*.png")) if fp_dir.is_dir() else []
    fn_files = sorted(fn_dir.glob("*.png")) if fn_dir.is_dir() else []

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Grid de Erros: FP e FN com Mapas Grad-CAM++",
                    f"{min(len(fp_files),8)} Falsos Positivos + {min(len(fn_files),8)} Falsos Negativos")

    if fp_files or fn_files:
        grid_buf = make_error_grid_png(fp_files, fn_files)
        slide.shapes.add_picture(grid_buf, Inches(0.15), Inches(1.35), Inches(13.0), Inches(6.0))
    else:
        add_text_box(slide,
                     "Imagens não encontradas em results/gradcam/.\n"
                     "Execute: python scripts/gradcam_error_analysis.py --checkpoint checkpoints/resnet50_full_finetune.pt",
                     Inches(1.0), Inches(3.0), Inches(11.0), Inches(2.0),
                     font_size=14, color=C_MID_GRAY, align=PP_ALIGN.CENTER)


def _challenge_card(slide, y: float, title: str, icon: str,
                    problem_lines: list[str], solution_lines: list[str],
                    status: str, color: RGBColor, status_color: RGBColor,
                    card_h: float = 1.9) -> None:
    """Renderiza um card de desafio com problema expandido, solução detalhada e status."""
    # Barra lateral colorida
    add_rect(slide, Inches(0.25), Inches(y), Inches(0.1), Inches(card_h - 0.08), fill=color)
    # Fundo do card
    add_rect(slide, Inches(0.38), Inches(y), Inches(12.6), Inches(card_h - 0.08),
             fill=RGBColor(0xF8, 0xF8, 0xF8), line=color, line_width_pt=0.6)
    # Título + ícone
    add_text_box(slide, f"{icon}  {title}",
                 Inches(0.5), Inches(y + 0.05), Inches(7.5), Inches(0.38),
                 font_size=13, bold=True, color=color)
    # Badge de status
    status_bg = RGBColor(0xD5, 0xFF, 0xD5) if "Resolvido" in status else RGBColor(0xFF, 0xF0, 0xD5)
    add_rect(slide, Inches(10.5), Inches(y + 0.08), Inches(2.4), Inches(0.28), fill=status_bg,
             line=status_color, line_width_pt=0.8)
    add_text_box(slide, status,
                 Inches(10.5), Inches(y + 0.08), Inches(2.4), Inches(0.28),
                 font_size=9, bold=True, color=status_color, align=PP_ALIGN.CENTER)
    # Seção Problema
    add_text_box(slide, "Problema:",
                 Inches(0.5), Inches(y + 0.48), Inches(1.5), Inches(0.25),
                 font_size=10, bold=True, color=C_MID_GRAY)
    add_bullet_list(slide, problem_lines,
                    Inches(0.5), Inches(y + 0.7), Inches(5.7), Inches(0.9),
                    font_size=10, color=C_DARK_GRAY, bullet_char="  ")
    # Separador vertical
    add_rect(slide, Inches(6.4), Inches(y + 0.45), Inches(0.03), Inches(card_h - 0.65),
             fill=RGBColor(0xCC, 0xCC, 0xCC))
    # Seção Solução
    add_text_box(slide, "Como foi resolvido:",
                 Inches(6.55), Inches(y + 0.48), Inches(3.0), Inches(0.25),
                 font_size=10, bold=True, color=C_GREEN)
    add_bullet_list(slide, solution_lines,
                    Inches(6.55), Inches(y + 0.7), Inches(6.2), Inches(0.9),
                    font_size=10, color=C_DARK_GRAY, bullet_char="✓  ")


def slide_challenges(prs: Presentation) -> None:
    """Slide 10a — 3 primeiros desafios com cards expandidos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Dificuldades e Soluções  (1/2)",
                    "Desafios técnicos encontrados e como cada um foi tratado")

    _challenge_card(
        slide, y=1.35,
        title="Dataset Desbalanceado  (NORMAL : PNEUMONIA = 1 : 2,9)",
        icon="⚖",
        problem_lines=[
            "Treino com 1.341 NORMAL vs 3.305 PNEUMONIA",
            "Sem correção, o modelo aprende a prever sempre PNEUMONIA",
            "e atinge ~74% de acurácia sem aprender nada útil",
        ],
        solution_lines=[
            "WeightedRandomSampler: cada batch tem proporção balanceada",
            "CrossEntropyLoss com pesos inversamente proporcionais",
            "  à frequência (w_NORMAL ≈ 1,96 × w_PNEUMONIA)",
            "Resultado: modelo aprende as duas classes corretamente",
        ],
        status="✓ Resolvido",
        color=C_ACCENT,
        status_color=C_GREEN,
    )

    _challenge_card(
        slide, y=3.4,
        title="Treinamento em CPU  (sem GPU disponível no ambiente)",
        icon="🖥",
        problem_lines=[
            "ResNet-50 em CPU: ~3–4 s/batch vs ~0,1 s/batch em GPU",
            "500 imgs de Grad-CAM no treino levaram ~5 minutos",
            "Inviável rodar muitas épocas ou datasets maiores",
        ],
        solution_lines=[
            "Estratégia head_only primeiro: backbone congelado, só FC",
            "  treina em minutos e serve como baseline sólido",
            "Número de épocas reduzido (5–10) com early stopping por AUC",
            "num_workers=0 para evitar erro de shared memory no WSL2",
        ],
        status="✓ Contornado",
        color=C_MED_BLUE,
        status_color=C_GREEN,
    )

    _challenge_card(
        slide, y=5.45,
        title="Overfitting no Full Finetune",
        icon="📉",
        problem_lines=[
            "AUC validação (0,9963) ≫ AUC teste (0,9693)",
            "Modelo memorizou o split de validação de 15%",
            "Split estático favorece vazamento de informação",
        ],
        solution_lines=[
            "Salvamento do melhor checkpoint por AUC de validação",
            "  (não a última época — evita piora por treino longo)",
            "XAI-guided fine-tuning mitigou parcialmente o efeito",
            "Próximo passo: validação cruzada k-fold + dropout",
        ],
        status="⚠ Parcial",
        color=C_RED,
        status_color=C_ACCENT,
    )

    # ── Slide 10b — 3 desafios restantes ──────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide2, "Dificuldades e Soluções  (2/2)",
                    "Desafios de implementação e limitações científicas em aberto")

    _challenge_card(
        slide2, y=1.35,
        title="Curva ROC Incompleta na Versão Inicial",
        icon="📊",
        problem_lines=[
            "Primeira versão reportava só o escalar AUC do JSON",
            "Não havia curva real com pontos (FPR, TPR) por threshold",
            "Figura era apenas um placeholder visual sem valor analítico",
        ],
        solution_lines=[
            "gradcam_error_analysis.py salva score_pneumonia por imagem",
            "generate_figures.py lê predictions[] e chama roc_curve(sklearn)",
            "Agora a curva usa os 624 scores reais do conjunto de teste",
            "AUC recalculado = 0,9692 (confirma o valor do JSON)",
        ],
        status="✓ Resolvido",
        color=C_GREEN,
        status_color=C_GREEN,
    )

    _challenge_card(
        slide2, y=3.4,
        title="Máscara Pulmonar Proxy Imprecisa",
        icon="🫁",
        problem_lines=[
            "Crop central de 60% não segmenta o pulmão de fato",
            "Regiões cardíacas e diafragma ficam dentro da máscara",
            "overlap_score superestima atenção 'focada' no pulmão real",
        ],
        solution_lines=[
            "Abordagem atual é conservadora: melhor que nenhuma máscara",
            "Opção 1: usar dataset Montgomery (anotações manuais de pulmão)",
            "Opção 2: U-Net pré-treinada para segmentação pulmonar",
            "Impacto prático: categorias focused/partial/distracted",
            "  são aproximadas mas suficientes para guiar o fine-tuning",
        ],
        status="⚠ Em aberto",
        color=RGBColor(0x7B, 0x68, 0xEE),
        status_color=RGBColor(0x7B, 0x68, 0xEE),
    )

    _challenge_card(
        slide2, y=5.45,
        title="Erro de Shared Memory no WSL2  (DataLoader multi-worker)",
        icon="💾",
        problem_lines=[
            "/dev/shm limitado no WSL2 — DataLoader com workers > 0 trava",
            "Erro: 'unable to allocate shared memory' no meio do treino",
            "Fase 1 do XAI fine-tuning completou, mas fase 2 falhou",
        ],
        solution_lines=[
            "num_workers=0 em todos os DataLoaders do xai_guided_finetune.py",
            "Solução simples: carregamento sequencial (um pouco mais lento)",
            "Alternativa futura: aumentar /dev/shm via wsl2.conf",
            "Fase 2 re-executada com sucesso após a correção",
        ],
        status="✓ Resolvido",
        color=RGBColor(0x80, 0x80, 0x00),
        status_color=C_GREEN,
    )


def slide_xai_finetune(prs: Presentation, m_xai: dict | None = None,
                       m_xai_v2: dict | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "XAI-Guided Fine-Tuning — Resultados Comparativos",
                    "Attention-Guided Sample Weighting iterado em V1 e V2")

    # Painel esquerdo — abordagem + tabela de 4 modelos
    add_rect(slide, Inches(0.3), Inches(1.35), Inches(7.8), Inches(5.65),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Abordagem: per-sample weight = 1 + α × overlap_score (α=2,0)",
                 Inches(0.5), Inches(1.4), Inches(7.5), Inches(0.38),
                 font_size=12, bold=True, color=C_DARK_BLUE)

    # Tabela 4 modelos
    cols_h = ["Modelo", "Acurácia", "AUC-ROC", "F1", "FP", "FN"]
    widths2 = [2.2, 1.3, 1.2, 1.2, 0.85, 0.85]
    xs2 = [0.35]
    for w in widths2[:-1]:
        xs2.append(xs2[-1] + w)

    for j, (c, w, x) in enumerate(zip(cols_h, widths2, xs2)):
        add_rect(slide, Inches(x), Inches(1.82), Inches(w), Inches(0.4), fill=C_DARK_BLUE)
        add_text_box(slide, c, Inches(x), Inches(1.82), Inches(w), Inches(0.4),
                     font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    t_xai   = m_xai["test"]   if m_xai   else None
    t_xai_v2 = m_xai_v2["test"] if m_xai_v2 else None
    rows2 = [
        ("Head Only",      "87.18%", "0.9300", "0.8942", "28", "52",  C_LIGHT_BLUE,         False),
        ("Full Finetune",  "86.54%", "0.9693", "0.9019", "80",  "4",  RGBColor(0xFF,0xF0,0xE0), False),
        ("XAI-Guided V1", f"{t_xai['accuracy_percent']:.2f}%" if t_xai else "—",
                          f"{t_xai['auc_roc']:.4f}"           if t_xai else "—",
                          f"{t_xai['f1']:.4f}"                if t_xai else "—",
                          str(t_xai['confusion_matrix'][0][1]) if t_xai else "—",
                          str(t_xai['confusion_matrix'][1][0]) if t_xai else "—",
                          RGBColor(0xD5, 0xFF, 0xD5), True),
        ("XAI-Guided V2", f"{t_xai_v2['accuracy_percent']:.2f}%" if t_xai_v2 else "—",
                          f"{t_xai_v2['auc_roc']:.4f}"            if t_xai_v2 else "—",
                          f"{t_xai_v2['f1']:.4f}"                 if t_xai_v2 else "—",
                          str(t_xai_v2['confusion_matrix'][0][1]) if t_xai_v2 else "—",
                          str(t_xai_v2['confusion_matrix'][1][0]) if t_xai_v2 else "—",
                          RGBColor(0xE0, 0xFF, 0xE8), True),
    ]
    for i, (nm, acc, auc_, f1, fp, fn, bg, is_xai) in enumerate(rows2):
        row_y = 2.26 + i * 0.52
        row_vals = [nm, acc, auc_, f1, fp, fn]
        for j, (v, w, x) in enumerate(zip(row_vals, widths2, xs2)):
            add_rect(slide, Inches(x), Inches(row_y), Inches(w), Inches(0.5), fill=bg)
            bold_it = is_xai and j in (1, 2, 3)
            add_text_box(slide, v, Inches(x), Inches(row_y), Inches(w), Inches(0.5),
                         font_size=11, bold=bold_it, color=C_DARK_GRAY, align=PP_ALIGN.CENTER)

    # Análise abaixo da tabela
    add_text_box(slide, "Evolução XAI-guided V2 (val):",
                 Inches(0.5), Inches(4.44), Inches(7.4), Inches(0.32),
                 font_size=11, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "Epoch 1 → val_acc=95.14%  AUC=0.9962  F1=0.9656",
        "Epoch 2 → val_acc=96.04%  AUC=0.9964  F1=0.9721",
        "Epoch 3 → val_acc=96.68%  AUC=0.9963  F1=0.9767",
    ], Inches(0.5), Inches(4.78), Inches(7.4), Inches(1.2), font_size=11)

    # Painel direito — análise
    add_rect(slide, Inches(8.35), Inches(1.35), Inches(4.95), Inches(5.65),
             fill=RGBColor(0xF0, 0xF0, 0xF0), line=C_ACCENT, line_width_pt=1.0)
    add_text_box(slide, "Análise & Tradeoffs",
                 Inches(8.55), Inches(1.4), Inches(4.6), Inches(0.4),
                 font_size=14, bold=True, color=C_ACCENT)

    fp_v1 = t_xai["confusion_matrix"][0][1]   if t_xai   else 38
    fn_v1 = t_xai["confusion_matrix"][1][0]   if t_xai   else 18
    fp_v2 = t_xai_v2["confusion_matrix"][0][1] if t_xai_v2 else 53
    fn_v2 = t_xai_v2["confusion_matrix"][1][0] if t_xai_v2 else  9

    add_bullet_list(slide, [
        f"V1 vs Full FT:  FP {80}→{fp_v1} (−{80-fp_v1}✓)  FN {4}→{fn_v1}",
        f"V2 vs V1:       FP {fp_v1}→{fp_v2}  FN {fn_v1}→{fn_v2} (−{fn_v1-fn_v2}✓)",
        "",
        "V1 = melhor equilíbrio geral",
        "  → +4,49% acurácia vs full_finetune",
        "  → +0,0281 F1  (mais balanceado)",
        "",
        "V2 = mais sensível (menos FN)",
        "  → custo: mais FP (53 vs 38)",
        "",
        "Ambos superam full_finetune",
        "em acurácia, F1 e erro total",
    ], Inches(8.55), Inches(1.88), Inches(4.6), Inches(4.6), font_size=11)


def slide_attention_comparison(prs: Presentation) -> None:
    """Slide com grade comparativa de FPs corrigidos e distribuição de overlap_score."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Comparação de Atenção Grad-CAM++ — Três Modelos",
                    "FPs corrigidos pelo XAI-guided e evolução do overlap_score")

    grade_pdf = FIGURES_DIR / "comparison_gradcam_3models.pdf"
    overlap_pdf = FIGURES_DIR / "overlap_score_distribution.pdf"
    evolution_pdf = FIGURES_DIR / "metrics_evolution.pdf"

    placed_any = False

    # Preferimos a figura de evolução de métricas (mais legível em slide)
    if evolution_pdf.is_file():
        import subprocess
        try:
            png_path = FIGURES_DIR / "_metrics_evolution_slide.png"
            if not png_path.is_file():
                subprocess.run(
                    ["pdftoppm", "-r", "150", "-png", "-singlefile",
                     str(evolution_pdf), str(FIGURES_DIR / "_metrics_evolution_slide")],
                    check=True, capture_output=True
                )
            if png_path.is_file():
                slide.shapes.add_picture(str(png_path),
                                         Inches(0.2), Inches(1.35), Inches(13.0), Inches(4.0))
                placed_any = True
        except Exception:
            pass

    if not placed_any:
        # fallback textual
        add_rect(slide, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                 fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
        add_text_box(slide, "Evolução de FP/FN por estratégia",
                     Inches(0.5), Inches(1.5), Inches(12.0), Inches(0.4),
                     font_size=14, bold=True, color=C_DARK_BLUE)

    # Linha de resumo embaixo
    add_rect(slide, Inches(0.3), Inches(5.55), Inches(12.7), Inches(1.5),
             fill=RGBColor(0xF0, 0xF0, 0xF0), line=C_ACCENT, line_width_pt=0.5)
    add_bullet_list(slide, [
        "Full Finetune: FP=80  FN=4   (máx recall, muitos alarmes falsos)",
        "XAI-Guided V1: FP=38  FN=18  (melhor equilíbrio — −53% FP, acurácia 91%)",
        "XAI-Guided V2: FP=53  FN=9   (mais sensível — menos FN, porém mais FP que V1)",
    ], Inches(0.5), Inches(5.6), Inches(12.4), Inches(1.4), font_size=12)


def slide_timeline(prs: Presentation, timeline_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Cronograma e Status do Projeto",
                    "7 semanas — abr/2026 a mai/2026")

    slide.shapes.add_picture(timeline_buf, Inches(0.3), Inches(1.35), Inches(12.7), Inches(3.0))

    # Resumo
    add_rect(slide, Inches(0.3), Inches(4.45), Inches(12.7), Inches(2.7),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Resumo de entregas",
                 Inches(0.5), Inches(4.5), Inches(12.0), Inches(0.4),
                 font_size=14, bold=True, color=C_DARK_BLUE)
    add_bullet_list(slide, [
        "✓  S1–S2: Ambiente Python, dataset, pipeline de dados completo",
        "✓  S3: ResNet-50 head_only e full_finetune treinados e avaliados",
        "✓  S4: Grad-CAM++ em lote (624 imagens), análise FP/FN, overlap_score",
        "✓  S5: Artigo LaTeX (abntex2) + figuras PDF + PowerPoint",
        "✓  S6: XAI-guided V1 (FP 80→38) e V2 (FN 18→9) concluídos",
        "✓  S7: Análise comparativa final, figuras e conclusões documentadas",
    ], Inches(0.5), Inches(4.9), Inches(12.4), Inches(2.1), font_size=13)


def slide_conclusion(prs: Presentation, m_full: dict,
                     m_xai: dict | None = None,
                     m_xai_v2: dict | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header_bar(slide, "Conclusão — Análise Comparativa Final",
                    "Quatro estratégias | transfer learning + explicabilidade + refinamento iterativo")

    # Painel esquerdo — resultados definitivos
    add_rect(slide, Inches(0.3), Inches(1.35), Inches(7.1), Inches(5.7),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.5)
    add_text_box(slide, "Resultados definitivos (conjunto de teste, 624 imgs)",
                 Inches(0.5), Inches(1.4), Inches(6.8), Inches(0.4),
                 font_size=13, bold=True, color=C_DARK_BLUE)

    xai_acc  = m_xai["test"]["accuracy_percent"]  if m_xai   else 91.03
    xai_f1   = m_xai["test"]["f1"]                if m_xai   else 0.9300
    xai_fp   = m_xai["test"]["confusion_matrix"][0][1] if m_xai else 38
    v2_fp    = m_xai_v2["test"]["confusion_matrix"][0][1] if m_xai_v2 else 53
    v2_fn    = m_xai_v2["test"]["confusion_matrix"][1][0] if m_xai_v2 else 9

    add_bullet_list(slide, [
        "Head Only:        acc=87,18%  AUC=0,9300  F1=0,8942  FP=28  FN=52",
        "Full Finetune:    acc=86,54%  AUC=0,9693  F1=0,9019  FP=80  FN=4",
        f"XAI-Guided V1 ★: acc={xai_acc:.2f}%  F1={xai_f1:.4f}  FP={xai_fp}  FN=18",
        f"XAI-Guided V2:   acc=90,06%  F1=0,9248  FP={v2_fp}  FN={v2_fn}",
        "",
        "■  XAI-Guided V1 = melhor equilíbrio geral",
        f"   +4,49% acurácia vs full_finetune",
        "   −52,5% FP (80→38) com acurácia e F1 superiores",
        "",
        "■  XAI-Guided V2 = mais sensível",
        f"   FN caem de 18 → {v2_fn}, mas FP sobem para {v2_fp}",
        "   Útil quando custo de FN > custo de FP",
    ], Inches(0.5), Inches(1.85), Inches(6.8), Inches(5.0), font_size=11)

    # Painel direito — lições e próximos passos
    add_rect(slide, Inches(7.65), Inches(1.35), Inches(5.6), Inches(2.8),
             fill=RGBColor(0xF0, 0xF0, 0xF0), line=C_ACCENT, line_width_pt=0.5)
    add_text_box(slide, "Lições aprendidas",
                 Inches(7.85), Inches(1.4), Inches(5.2), Inches(0.4),
                 font_size=13, bold=True, color=C_ACCENT)
    add_bullet_list(slide, [
        "XAI não é só diagnóstico — pode guiar o treino",
        "Grad-CAM++ + overlap_score = sinal de qualidade por amostra",
        "Iterações adicionais de XAI não eliminam tradeoff FP/FN",
        "Pipeline reproduzível: checkpoint → CAM → pesos → fine-tune",
    ], Inches(7.85), Inches(1.85), Inches(5.2), Inches(2.15), font_size=11)

    add_rect(slide, Inches(7.65), Inches(4.35), Inches(5.6), Inches(2.7),
             fill=C_LIGHT_BLUE, line=C_GREEN, line_width_pt=0.8)
    add_text_box(slide, "Trabalhos Futuros",
                 Inches(7.85), Inches(4.4), Inches(5.2), Inches(0.4),
                 font_size=13, bold=True, color=C_GREEN)
    add_bullet_list(slide, [
        "Segmentação pulmonar real (substituir proxy)",
        "Calibração de threshold de decisão pós-treino",
        "Validação cruzada k-fold",
        "Discussão clínica qualitativa com especialista",
        "Extensão para classificação multiclasse",
    ], Inches(7.85), Inches(4.85), Inches(5.2), Inches(2.0), font_size=11)


# ──────────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────────

def load_json(p: Path) -> dict:
    return json.loads(p.read_text(encoding="utf-8"))


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path,
                        default=ROOT / "apresentacao.pptx",
                        help="Caminho de saída do .pptx")
    args = parser.parse_args()

    # ── Carrega dados ──────────────────────────────────────────────────────────
    metrics_head_path  = RESULTS_DIR / "metrics_head_only.json"
    metrics_full_path  = RESULTS_DIR / "metrics_full_finetune.json"
    metrics_xai_path   = RESULTS_DIR / "metrics_xai_guided.json"
    metrics_xai_v2_path = RESULTS_DIR / "metrics_xai_guided_v2.json"
    error_json_path    = RESULTS_DIR / "error_analysis.json"

    if not metrics_head_path.is_file() or not metrics_full_path.is_file():
        print("ERRO: arquivos de métricas não encontrados em results/", file=sys.stderr)
        return 2

    m_head   = load_json(metrics_head_path)
    m_full   = load_json(metrics_full_path)
    m_xai    = load_json(metrics_xai_path)    if metrics_xai_path.is_file()    else None
    m_xai_v2 = load_json(metrics_xai_v2_path) if metrics_xai_v2_path.is_file() else None
    error_json = load_json(error_json_path) if error_json_path.is_file() else {}
    predictions = error_json.get("predictions", [])

    # ── Gera figuras em memória ────────────────────────────────────────────────
    print("Gerando figuras...")
    cm = np.array(m_full["test"]["confusion_matrix"], dtype=int)
    cm_buf   = make_confusion_matrix_png(cm, "Matriz de Confusão — Full Finetune (teste)")
    roc_buf  = make_roc_png(predictions, "Curva ROC — Full Finetune") if predictions else None
    bar_buf  = make_error_bar_png(
        error_json.get("error_category_counts", {}),
        "Categorias de Atenção nos Erros"
    )
    cmp_buf   = make_metrics_comparison_png(m_head, m_full)
    arch_buf  = make_architecture_png()
    xai_buf   = make_xai_concept_png()
    tl_buf    = make_timeline_png()

    compare_img = ROOT / "results" / "gradcam_compare" / "person1_bacteria_1_gradcam++_A_vs_B.png"
    grid_pdf    = FIGURES_DIR / "error_grid_4x4_full_finetune.pdf"

    # ── Cria apresentação ──────────────────────────────────────────────────────
    print("Montando slides...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_motivation(prs)
    slide_dataset(prs)
    slide_strategies(prs, m_head, m_full)        # slide 4 — head_only vs full_finetune
    slide_architecture(prs, arch_buf)             # slide 5 — diagrama ResNet-50
    slide_results_table(prs, m_head, m_full)      # slide 6 — tabela métricas head+full
    slide_confusion_roc(prs, cm_buf, roc_buf or cm_buf, m_full)   # slide 7
    slide_gradcam_intro(prs, compare_img, xai_buf)                # slide 8
    slide_error_analysis(prs, bar_buf, error_json)                # slide 9
    slide_error_grid(prs, grid_pdf)                               # slide 10
    slide_challenges(prs)                                         # slides 11-12
    slide_xai_finetune(prs, m_xai, m_xai_v2)                     # slide 13 — XAI V1+V2
    slide_attention_comparison(prs)                               # slide 14 — evolução atenção
    slide_timeline(prs, tl_buf)                                   # slide 15
    slide_conclusion(prs, m_full, m_xai, m_xai_v2)               # slide 16 — conclusão final

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print(f"\nApresentação salva em: {args.out}")
    print(f"Total de slides: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
