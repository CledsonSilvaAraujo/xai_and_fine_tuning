#!/usr/bin/env python3
"""
Gera a APRESENTAÇÃO FINAL (3ª parte) do projeto XAI-Guided Fine-Tuning.

Narrativa: fecha o ciclo Proposta (Apres. 1) → Execução (Apres. 2) → Conclusão (Apres. 3).
Inclui todos os 4 modelos: head_only, full_finetune, XAI-guided V1, XAI-guided V2.

Uso:
  python scripts/create_presentation_final.py
  python scripts/create_presentation_final.py --out apresentacao_final.pptx
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns
from PIL import Image
from sklearn.metrics import auc, roc_curve

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT        = Path(__file__).resolve().parents[1]
RESULTS_DIR = ROOT / "results"
FIGURES_DIR = ROOT / "figures"

# ─── Paleta ──────────────────────────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x10, 0x3A, 0x6E)
C_MED_BLUE   = RGBColor(0x1A, 0x6C, 0xB0)
C_LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
C_ACCENT     = RGBColor(0xE8, 0x6A, 0x2B)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
C_MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
C_GREEN      = RGBColor(0x2C, 0xA0, 0x2C)
C_RED        = RGBColor(0xD6, 0x27, 0x28)
C_YELLOW     = RGBColor(0xFF, 0xBF, 0x00)
C_PURPLE     = RGBColor(0x7B, 0x68, 0xEE)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill=None, line=None, line_width_pt=0.5):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width_pt)
    return shape


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, italic=False,
             color=C_DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullets(slide, items, left, top, width, height,
                size=14, color=C_DARK_GRAY, bullet="•  "):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet + item if item.strip() else ""
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def make_header(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill=C_DARK_BLUE)
    add_text(slide, title, Inches(0.35), Inches(0.05), Inches(12.6), Inches(0.75),
             size=22, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.8), Inches(12.6), Inches(0.45),
                 size=13, italic=True, color=RGBColor(0xC0, 0xD8, 0xF0))


def fig_to_bytes(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def load_json(p: Path) -> dict:
    return json.loads(p.read_text(encoding="utf-8"))


# ─── Figuras ─────────────────────────────────────────────────────────────────

def make_pipeline_png() -> io.BytesIO:
    """Diagrama visual do pipeline completo dos 4 modelos."""
    fig, ax = plt.subplots(figsize=(12, 2.6))
    ax.axis("off")
    fig.patch.set_facecolor("#f5f8fc")

    steps = [
        ("Dados\nNIH Chest\nX-Ray", "#4472C4", "white"),
        ("Head Only\n5 épocas\n4K params", "#7030A0", "white"),
        ("Full\nFinetune\n10 épocas\n22M params", "#ED7D31", "white"),
        ("XAI-Guided\nV1\n3 épocas\nLR=1e-5", "#2CA02C", "white"),
        ("XAI-Guided\nV2\n3 épocas\nLR=5e-6", "#1A6CB0", "white"),
        ("Análise\nComparativa\nFinal", "#D62728", "white"),
    ]
    w, h, gap = 1.7, 1.8, 0.18
    total = len(steps) * w + (len(steps) - 1) * gap
    start_x = (12 - total) / 2

    for i, (label, color, tc) in enumerate(steps):
        x = start_x + i * (w + gap)
        rect = mpatches.FancyBboxPatch((x, 0.3), w, h,
                                        boxstyle="round,pad=0.07",
                                        facecolor=color, edgecolor="white", linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + w / 2, 0.3 + h / 2, label,
                ha="center", va="center", fontsize=8.5,
                fontweight="bold", color=tc, linespacing=1.4)
        if i < len(steps) - 1:
            ax.annotate("", xy=(x + w + gap, 0.3 + h / 2),
                        xytext=(x + w, 0.3 + h / 2),
                        arrowprops=dict(arrowstyle="->", color="#555", lw=2))

    # Labels embaixo indicando "Pres. 1", "Pres. 2", "Pres. 3"
    brackets = [
        (start_x, start_x + w + gap + w * 0.3, "Apresentação 1\n(Proposta)", "#4472C4"),
        (start_x + w + gap + w * 0.4, start_x + 3 * (w + gap) + w * 0.5,
         "Apresentação 2\n(Entrega Intermediária)", "#ED7D31"),
        (start_x + 3 * (w + gap) + w * 0.4, start_x + total - 0.1,
         "Apresentação 3\n(Conclusão Final — esta)", "#D62728"),
    ]
    for x0, x1, lbl, clr in brackets:
        mid = (x0 + x1) / 2
        ax.annotate("", xy=(x1, 0.12), xytext=(x0, 0.12),
                    arrowprops=dict(arrowstyle="-", color=clr, lw=1.5))
        ax.text(mid, -0.02, lbl, ha="center", va="top", fontsize=7.5, color=clr, fontweight="bold")

    ax.set_xlim(0, 12)
    ax.set_ylim(-0.25, 2.5)
    plt.tight_layout(pad=0.1)
    return fig_to_bytes(fig)


def make_4models_bar_png(models_data: list[tuple]) -> io.BytesIO:
    """Gráfico de barras agrupadas: FP e FN dos 4 modelos."""
    names = [m[0] for m in models_data]
    fps   = [m[1] for m in models_data]
    fns   = [m[2] for m in models_data]
    colors_fp = ["#9467BD", "#1A6CB0", "#2CA02C", "#17BECF"]
    colors_fn = ["#C49AD8", "#8EC8F5", "#90E890", "#9EDAE5"]

    x = np.arange(len(names))
    w = 0.35
    fig, ax = plt.subplots(figsize=(9, 4.5))
    bars1 = ax.bar(x - w/2, fps, w, label="FP (alarme falso)", color=colors_fp, edgecolor="white")
    bars2 = ax.bar(x + w/2, fns, w, label="FN (pneumonia perdida)", color=colors_fn,
                   edgecolor="white", hatch="//")

    for bar in list(bars1) + list(bars2):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.8,
                str(int(bar.get_height())), ha="center", va="bottom",
                fontweight="bold", fontsize=12)

    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=11)
    ax.set_ylabel("Número de erros", fontsize=11)
    ax.set_title("FP e FN por Modelo — Conjunto de Teste (624 imgs)", fontsize=13, fontweight="bold")
    ax.legend(fontsize=10)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def make_metrics_evolution_png(models_data: list[tuple]) -> io.BytesIO:
    """Gráfico de linhas: evolução de AUC, F1 e Acurácia."""
    names  = [m[0] for m in models_data]
    aucs   = [m[3] for m in models_data]
    f1s    = [m[4] for m in models_data]
    accs   = [m[5]/100 for m in models_data]

    x = list(range(len(names)))
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.set_facecolor("#f9f9f9")
    fig.patch.set_facecolor("white")

    ax.plot(x, aucs, "o-",  color="#1A6CB0", lw=2.5, ms=9, label="AUC-ROC")
    ax.plot(x, f1s,  "s--", color="#E86A2B", lw=2.5, ms=9, label="F1")
    ax.plot(x, accs, "^:",  color="#2CA02C", lw=2.5, ms=9, label="Acurácia")

    for i, (a, f, ac) in enumerate(zip(aucs, f1s, accs)):
        ax.annotate(f"{a:.4f}", (x[i], a), textcoords="offset points",
                    xytext=(0, 9), ha="center", fontsize=8.5, color="#1A6CB0", fontweight="bold")
        ax.annotate(f"{f:.4f}", (x[i], f), textcoords="offset points",
                    xytext=(0, -16), ha="center", fontsize=8.5, color="#E86A2B", fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=11)
    ax.set_ylim(0.83, 1.01)
    ax.set_title("Evolução de AUC-ROC, F1 e Acurácia", fontsize=13, fontweight="bold")
    ax.legend(fontsize=10, loc="lower right")
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def make_overlap_comparison_png() -> io.BytesIO:
    """Boxplot + jitter do overlap_score nos erros dos 3 modelos pós full_finetune."""
    models_overlap = [
        ("Full Finetune",   RESULTS_DIR / "error_analysis.json",        "#1A6CB0"),
        ("XAI-Guided V1",  RESULTS_DIR / "error_analysis_xai.json",     "#2CA02C"),
        ("XAI-Guided V2",  RESULTS_DIR / "error_analysis_xai_v2.json",  "#17BECF"),
    ]

    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.set_facecolor("#f9f9f9")
    rng = np.random.default_rng(42)

    positions = []
    groups = []
    for i, (label, path, clr) in enumerate(models_overlap):
        if not path.is_file():
            continue
        data = json.loads(path.read_text())
        scores = [e["overlap_score"] for e in data.get("errors", []) if "overlap_score" in e]
        if not scores:
            continue
        positions.append(i + 1)
        groups.append((label, scores, clr))

    if not groups:
        ax.text(0.5, 0.5, "Dados não disponíveis", ha="center", va="center",
                transform=ax.transAxes, fontsize=14)
        return fig_to_bytes(fig)

    bp = ax.boxplot([g[1] for g in groups], positions=[p for p in range(1, len(groups)+1)],
                    widths=0.4, patch_artist=True, notch=False,
                    medianprops=dict(color="black", linewidth=2))
    for patch, (_, _, clr) in zip(bp["boxes"], groups):
        patch.set_facecolor(clr)
        patch.set_alpha(0.55)

    for pos, (lbl, scores, clr) in enumerate(groups, start=1):
        jitter = rng.uniform(-0.13, 0.13, size=len(scores))
        ax.scatter(pos + jitter, scores, color=clr, alpha=0.5, s=16, zorder=3)

    ax.axhline(0.20, linestyle="--", color="#888", lw=1, label="focused <0.20")
    ax.axhline(0.50, linestyle=":",  color="#555", lw=1, label="distracted >0.50")

    ax.set_xticks(list(range(1, len(groups)+1)))
    ax.set_xticklabels([g[0] for g in groups], fontsize=11)
    ax.set_ylabel("overlap_score (fração fora da máscara proxy)", fontsize=10)
    ax.set_title("Distribuição do overlap_score nos Erros", fontsize=13, fontweight="bold")
    ax.legend(fontsize=9, loc="upper right")

    for pos, (lbl, scores, clr) in enumerate(groups, start=1):
        ax.text(pos, ax.get_ylim()[0] - 0.01,
                f"n={len(scores)}\nmed={np.median(scores):.3f}",
                ha="center", va="top", fontsize=8.5, color=clr, fontweight="bold")

    fig.tight_layout()
    return fig_to_bytes(fig)


def make_tradeoff_diagram_png(models_data: list[tuple]) -> io.BytesIO:
    """Scatter plot FP × FN com os 4 modelos posicionados."""
    fig, ax = plt.subplots(figsize=(7, 5))
    ax.set_facecolor("#f9f9f9")

    colors = ["#7B68EE", "#1A6CB0", "#2CA02C", "#17BECF"]
    markers = ["D", "s", "★", "^"]

    for i, (nm, fp, fn, auc_, f1, acc) in enumerate(models_data):
        ax.scatter(fp, fn, color=colors[i], s=220, zorder=5, label=f"{nm}  (acc={acc:.1f}%)")
        ax.annotate(nm, (fp, fn), textcoords="offset points",
                    xytext=(8, 5), fontsize=9.5, color=colors[i], fontweight="bold")

    ax.set_xlabel("Falsos Positivos (FP) — alarmes falsos", fontsize=11)
    ax.set_ylabel("Falsos Negativos (FN) — pneumonias não detectadas", fontsize=11)
    ax.set_title("Tradeoff FP × FN: escolha o modelo conforme o contexto clínico",
                 fontsize=12, fontweight="bold")
    ax.legend(fontsize=9, loc="upper right")
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.xaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)

    # Anotações de contexto
    ax.annotate("← Menos FP\nmenos alarmes falsos",
                xy=(35, 52), fontsize=8, color="#2CA02C",
                ha="center", style="italic")
    ax.annotate("↓ Menos FN\nmenos pneumonias perdidas",
                xy=(85, 7), fontsize=8, color="#D62728",
                ha="center", style="italic")

    fig.tight_layout()
    return fig_to_bytes(fig)


# ─── Slides ──────────────────────────────────────────────────────────────────

def slide_capa(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Fundo superior
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_DARK_BLUE)
    # Faixa central
    add_rect(slide, Inches(0), Inches(2.6), SLIDE_W, Inches(2.5),
             fill=RGBColor(0x1A, 0x6C, 0xB0))

    add_text(slide, "Explicabilidade em Visão Computacional",
             Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.7),
             size=20, bold=False, italic=True, color=RGBColor(0xC0, 0xD8, 0xF0),
             align=PP_ALIGN.CENTER)
    add_text(slide, "para Diagnóstico Médico",
             Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
             size=20, bold=False, italic=True, color=RGBColor(0xC0, 0xD8, 0xF0),
             align=PP_ALIGN.CENTER)

    add_text(slide, "Apresentação Final",
             Inches(0.5), Inches(2.65), Inches(12.3), Inches(0.65),
             size=28, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, "Análise Comparativa Final — 4 Modelos",
             Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.6),
             size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "XAI-Guided Fine-Tuning: Proposta → Execução → Conclusão",
             Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.5),
             size=16, italic=True, color=RGBColor(0xE0, 0xF0, 0xFF),
             align=PP_ALIGN.CENTER)

    add_text(slide,
             "Cledson da Silva Araujo  ·  Visão Computacional  ·  UFF — 2026",
             Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
             size=14, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)
    add_text(slide,
             "ResNet-50  |  Grad-CAM++  |  Chest X-Ray NIH",
             Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
             size=13, italic=True, color=RGBColor(0x90, 0xB8, 0xD8),
             align=PP_ALIGN.CENTER)

    # Badges (3 apresentações)
    badge_labels = ["Apres. 1\nProposta", "Apres. 2\nEntrega\nIntermediária", "Apres. 3\nConclusão\nFinal"]
    badge_colors = [RGBColor(0x44, 0x72, 0xC4), RGBColor(0xED, 0x7D, 0x31), RGBColor(0xD6, 0x27, 0x28)]
    xs = [1.3, 5.6, 9.9]
    for lbl, clr, x in zip(badge_labels, badge_colors, xs):
        add_rect(slide, Inches(x), Inches(6.5), Inches(2.4), Inches(0.75), fill=clr)
        add_text(slide, lbl, Inches(x), Inches(6.5), Inches(2.4), Inches(0.75),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_ciclo_completo(prs: Presentation, pipeline_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Do Problema à Conclusão: O Ciclo Completo",
                "Três apresentações — um projeto coeso de XAI aplicado a diagnóstico médico")

    slide.shapes.add_picture(pipeline_buf, Inches(0.3), Inches(1.35), Inches(12.7), Inches(2.6))

    # Três cards na base
    cards = [
        ("Apresentação 1 — Proposta",
         ["Problema: CNN como caixa-preta em diagnóstico",
          "Pergunta: quais pixels causam o erro?",
          "Pipeline XAI-Guided (4 passos)",
          "Diferencial: sem anotações manuais",
          "Cronograma 2 fases / 7 semanas"],
         RGBColor(0x44, 0x72, 0xC4), 0.3),
        ("Apresentação 2 — Entrega Intermediária",
         ["Dataset, preprocessing, estratégias",
          "Head Only vs Full Finetune treinados",
          "Grad-CAM++ em 624 imagens",
          "XAI-Guided V1 implementado (FP 80→38)",
          "Artigo LaTeX + figuras entregues"],
         RGBColor(0xED, 0x7D, 0x31), 4.6),
        ("Apresentação 3 — Conclusão (esta)",
         ["XAI-Guided V2: FN 18→9",
          "Comparação completa: 4 modelos",
          "Análise de atenção Grad-CAM++ evoluída",
          "Tradeoff FP/FN e implicações clínicas",
          "Lições aprendidas + trabalhos futuros"],
         RGBColor(0xD6, 0x27, 0x28), 8.9),
    ]
    for title, bullets, clr, x in cards:
        add_rect(slide, Inches(x), Inches(4.15), Inches(4.2), Inches(3.0),
                 fill=RGBColor(0xF5, 0xF5, 0xF5), line=clr, line_width_pt=1.5)
        add_text(slide, title, Inches(x + 0.1), Inches(4.2), Inches(4.0), Inches(0.45),
                 size=11, bold=True, color=clr)
        add_bullets(slide, bullets, Inches(x + 0.1), Inches(4.65),
                    Inches(4.0), Inches(2.4), size=10)


def slide_promessa_vs_entregue(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "O que foi Prometido vs O que foi Entregue",
                "Proposta (Apresentação 1) confrontada com os resultados finais")

    rows = [
        ("Treinar ResNet-50 com transfer learning",             "✓ head_only e full_finetune treinados",           True),
        ("Gerar mapas Grad-CAM++ sobre os erros",              "✓ 624 imgs processadas; overlays salvos",          True),
        ("Análise sistemática de erros (categorização)",       "✓ focused / partial / distracted por overlap_score", True),
        ("Fine-tuning guiado por XAI (custo ponderado)",       "✓ V1: FP 80→38 (−52,5%), acc 86%→91%",           True),
        ("Sem anotações manuais — sinal automático",           "✓ overlap_score via máscara proxy 60%",           True),
        ("Segunda iteração XAI-Guided",                        "✓ V2: FN 18→9 (ao custo de FP=53)",              True),
        ("Métricas: AUC, F1, Precisão, Recall, AOPC",         "⚠ AOPC não implementado — AUC/F1/FP/FN realizados", False),
        ("Segmentação pulmonar real (U-Net / Montgomery)",     "⚠ Máscara proxy — segmentação real: trabalho futuro", False),
        ("Validação cruzada k-fold",                           "⚠ Split estático — k-fold: trabalho futuro",      False),
    ]

    # Cabeçalho da tabela
    hdrs = ["Compromisso (Apres. 1)", "Resultado (Apres. 3)", "Status"]
    ws = [5.5, 5.8, 1.7]
    xs = [0.25, 5.8, 11.65]
    for j, (h, w, x) in enumerate(zip(hdrs, ws, xs)):
        add_rect(slide, Inches(x), Inches(1.35), Inches(w), Inches(0.42), fill=C_DARK_BLUE)
        add_text(slide, h, Inches(x), Inches(1.35), Inches(w), Inches(0.42),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for i, (comp, result, ok) in enumerate(rows):
        bg = RGBColor(0xE8, 0xFF, 0xE8) if ok else RGBColor(0xFF, 0xF4, 0xE0)
        y = 1.8 + i * 0.54
        for j, (v, w, x) in enumerate(zip([comp, result, "✓ OK" if ok else "⚠ Parcial"], ws, xs)):
            add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(0.52), fill=bg,
                     line=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.3)
            clr = C_GREEN if (j == 2 and ok) else (C_ACCENT if (j == 2 and not ok) else C_DARK_GRAY)
            add_text(slide, v, Inches(x + 0.05), Inches(y), Inches(w - 0.1), Inches(0.52),
                     size=9.5, color=clr, bold=(j == 2), align=PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT)


def slide_pipeline_final(prs: Presentation, pipeline_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Pipeline Final: Quatro Modelos, Uma Arquitetura",
                "ResNet-50 com estratégias progressivas de fine-tuning e guia XAI")

    slide.shapes.add_picture(pipeline_buf, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.5))

    # Painel de estratégias
    strategies = [
        ("HEAD ONLY", "4.098 param\n5 épocas\nLR=1e-3\nBaseline rápido",
         "#7B68EE", C_DARK_BLUE),
        ("FULL FINETUNE", "~22M param\n10 épocas\nLR Layer3/4=1e-5\nAlto recall",
         "#1A6CB0", C_DARK_BLUE),
        ("XAI-GUIDED V1", "Parte do FF\n3 épocas\nLR=1e-5\nPeso por overlap",
         "#2CA02C", C_DARK_BLUE),
        ("XAI-GUIDED V2", "Parte do V1\n3 épocas\nLR=5e-6\nMais sensível",
         "#17BECF", C_DARK_BLUE),
    ]
    for i, (title, body, bg_hex, fg) in enumerate(strategies):
        x = 0.3 + i * 3.27
        bg = RGBColor(int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16))
        add_rect(slide, Inches(x), Inches(4.1), Inches(3.1), Inches(3.2), fill=bg)
        add_text(slide, title, Inches(x), Inches(4.15), Inches(3.1), Inches(0.5),
                 size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, body, Inches(x + 0.1), Inches(4.7), Inches(2.9), Inches(2.5),
                 size=11, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_resultados_tabela(prs: Presentation, models_data: list[tuple]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Resultados Comparativos — Quatro Estratégias",
                "Conjunto de teste: 624 imagens (234 NORMAL / 390 PNEUMONIA)")

    # Tabela principal
    cols = ["Modelo", "Acurácia", "Erro", "AUC-ROC", "F1", "Precisão", "Recall", "FP", "FN"]
    widths = [2.2, 1.25, 1.0, 1.2, 1.1, 1.2, 1.1, 0.85, 0.85]
    xs = [0.25]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)

    bg_colors = [
        RGBColor(0x7B, 0x68, 0xEE),   # head_only — roxo
        RGBColor(0x1A, 0x6C, 0xB0),   # full_finetune — azul
        RGBColor(0x2C, 0xA0, 0x2C),   # xai v1 — verde
        RGBColor(0x17, 0xBE, 0xCF),   # xai v2 — ciano
    ]

    # Cabeçalho
    for j, (c, w, x) in enumerate(zip(cols, widths, xs)):
        add_rect(slide, Inches(x), Inches(1.38), Inches(w), Inches(0.5), fill=C_DARK_BLUE)
        add_text(slide, c, Inches(x), Inches(1.38), Inches(w), Inches(0.5),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    row_h = 0.88
    for i, (nm, fp, fn, auc_, f1, acc) in enumerate(models_data):
        prec_map = {"Head Only": 0.9235, "Full Finetune": 0.8283,
                    "XAI-Guided V1": 0.9073, "XAI-Guided V2": 0.8779}
        rec_map  = {"Head Only": 0.8667, "Full Finetune": 0.9897,
                    "XAI-Guided V1": 0.9538, "XAI-Guided V2": 0.9769}
        prec = prec_map.get(nm, 0)
        rec  = rec_map.get(nm, 0)
        err  = 100 - acc

        vals = [nm, f"{acc:.2f}%", f"{err:.2f}%", f"{auc_:.4f}",
                f"{f1:.4f}", f"{prec:.4f}", f"{rec:.4f}", str(fp), str(fn)]

        y = 1.9 + i * row_h
        is_best = (nm == "XAI-Guided V1")
        bg = RGBColor(0xD5, 0xFF, 0xD5) if is_best else (
             RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else C_WHITE)

        for j, (v, w, x) in enumerate(zip(vals, widths, xs)):
            cell_bg = bg_colors[i] if j == 0 else bg
            add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(row_h - 0.05),
                     fill=cell_bg, line=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.3)
            fc = C_WHITE if j == 0 else (C_GREEN if is_best else C_DARK_GRAY)
            add_text(slide, v, Inches(x), Inches(y + 0.2), Inches(w), Inches(row_h - 0.25),
                     size=12 if j == 0 else 11,
                     bold=(j == 0 or (is_best and j in (1, 3, 4))),
                     color=fc, align=PP_ALIGN.CENTER)

    # Legenda
    add_text(slide,
             "★ XAI-Guided V1 = melhor equilíbrio  |  V2 = mais sensível (menos FN)  "
             "|  Full Finetune = máx recall  |  Head Only = máx precisão",
             Inches(0.25), Inches(7.05), Inches(13.0), Inches(0.35),
             size=10, italic=True, color=C_MID_GRAY, align=PP_ALIGN.CENTER)


def slide_evolucao_metricas(prs: Presentation, evo_buf: io.BytesIO,
                             bar_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Evolução das Métricas ao Longo das Estratégias",
                "Cada modelo é uma iteração sobre o anterior — do baseline à conclusão")

    slide.shapes.add_picture(evo_buf, Inches(0.2), Inches(1.35), Inches(6.6), Inches(4.5))
    slide.shapes.add_picture(bar_buf, Inches(6.9), Inches(1.35), Inches(6.2), Inches(4.5))

    add_text(slide,
             "AUC e F1 sobem com XAI-Guided V1. V2 mantém AUC mas ajusta o tradeoff FP/FN.",
             Inches(0.3), Inches(6.05), Inches(12.7), Inches(0.35),
             size=11, italic=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_xai_metodo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "XAI-Guided Fine-Tuning: O Método em Detalhe",
                "Sinal de supervisão automático derivado da atenção Grad-CAM++")

    # 4 passos
    steps = [
        ("1", "Checkpoint\nFull Finetune", "Ponto de partida:\nmodelo já treinado\nem radiografias"),
        ("2", "Grad-CAM++\n500 imgs treino", "overlap_score por imagem:\n% ativação FORA\nda máscara pulmonar"),
        ("3", "Per-sample\nWeight", "w = 1 + 2 × overlap_score\nAmostras 'distraídas'\nrecebem peso maior"),
        ("4", "Fine-Tune\n3 épocas", "LR = 1e-5 (V1)\nLR = 5e-6 (V2)\nBackbone × 0,1"),
    ]
    colors = ["#1A6CB0", "#E86A2B", "#2CA02C", "#7B68EE"]
    for i, (num, title, body) in enumerate(steps):
        x = 0.3 + i * 3.27
        bg = RGBColor(int(colors[i][1:3], 16), int(colors[i][3:5], 16), int(colors[i][5:7], 16))
        add_rect(slide, Inches(x), Inches(1.4), Inches(3.0), Inches(3.5), fill=bg)
        add_rect(slide, Inches(x + 0.85), Inches(1.45), Inches(1.3), Inches(0.7),
                 fill=C_WHITE)
        add_text(slide, num, Inches(x + 0.85), Inches(1.45), Inches(1.3), Inches(0.7),
                 size=22, bold=True, color=bg, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(x + 0.1), Inches(2.2), Inches(2.8), Inches(0.8),
                 size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, body, Inches(x + 0.1), Inches(3.05), Inches(2.8), Inches(1.75),
                 size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

        if i < 3:
            add_text(slide, "→", Inches(x + 2.85), Inches(2.6), Inches(0.5), Inches(0.5),
                     size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Painel de resultados
    add_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(2.15),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.8)
    add_text(slide, "Resultado: o modelo é penalizado quando olha para o lugar errado",
             Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.4),
             size=13, bold=True, color=C_DARK_BLUE)
    add_bullets(slide, [
        "V1 (LR=1e-5): FP 80 → 38  (−52,5%)  |  FN 4 → 18  |  Acurácia 86,5% → 91,0%  |  F1 +0,028",
        "V2 (LR=5e-6, parte do V1): FP 38 → 53  |  FN 18 → 9 (−50%)  |  Acurácia 90,1%",
        "Princípio: Grad-CAM++ como sinal de qualidade por amostra — sem nenhuma anotação manual adicional",
    ], Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), size=11)


def slide_v2_deep_dive(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "XAI-Guided V2: A Segunda Iteração",
                "O que muda ao iterar o processo — e o que aprendemos com o resultado")

    # Configuração
    add_rect(slide, Inches(0.3), Inches(1.38), Inches(4.6), Inches(5.7),
             fill=RGBColor(0xE0, 0xEE, 0xFF), line=C_MED_BLUE, line_width_pt=0.8)
    add_text(slide, "Configuração V2", Inches(0.5), Inches(1.43), Inches(4.3), Inches(0.42),
             size=13, bold=True, color=C_DARK_BLUE)
    add_bullets(slide, [
        "Checkpoint: resnet50_xai_guided.pt (V1)",
        "LR = 5e-6  (÷2 em relação ao V1)",
        "Épocas: 3  |  α = 2,0",
        "500 imgs treino para CAM weights",
        "",
        "Overlap scores (treino):",
        "  3 focused  |  471 partial  |  26 distracted",
        "  média=0,369  mediana=0,366",
        "",
        "Evolução de validação:",
        "  Epoch 1 → acc=95,14%  AUC=0,9962  F1=0,9656",
        "  Epoch 2 → acc=96,04%  AUC=0,9964  F1=0,9721",
        "  Epoch 3 → acc=96,68%  AUC=0,9963  F1=0,9767",
    ], Inches(0.5), Inches(1.9), Inches(4.3), Inches(5.1), size=10.5)

    # Resultados V1 vs V2
    add_rect(slide, Inches(5.2), Inches(1.38), Inches(7.9), Inches(3.0),
             fill=RGBColor(0xF0, 0xF0, 0xF0), line=C_ACCENT, line_width_pt=0.8)
    add_text(slide, "V1 vs V2 — Conjunto de Teste",
             Inches(5.4), Inches(1.43), Inches(7.5), Inches(0.42),
             size=13, bold=True, color=C_ACCENT)

    comp_rows = [
        ("Métrica", "Full FT", "XAI V1", "XAI V2"),
        ("Acurácia",  "86,54%", "91,03% ↑", "90,06%"),
        ("AUC-ROC",   "0,9693", "0,9663",   "0,9655"),
        ("F1",        "0,9019", "0,9300 ↑", "0,9248"),
        ("FP",        "80",     "38 ↓",     "53"),
        ("FN",        "4",      "18",        "9 ↓"),
    ]
    cws = [2.0, 1.7, 1.85, 1.85]
    cxs = [5.25]
    for w in cws[:-1]:
        cxs.append(cxs[-1] + w)

    for r, row in enumerate(comp_rows):
        for c, (v, w, x) in enumerate(zip(row, cws, cxs)):
            bg = C_DARK_BLUE if r == 0 else (
                 RGBColor(0xD5, 0xFF, 0xD5) if "↑" in v or "↓" in v else
                 (RGBColor(0xF5, 0xF5, 0xF5) if r % 2 == 0 else C_WHITE))
            add_rect(slide, Inches(x), Inches(1.9 + r * 0.39), Inches(w), Inches(0.38),
                     fill=bg, line=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.3)
            fc = C_WHITE if r == 0 else (C_GREEN if "↑" in v else
                                          (C_RED if "↓" in v else C_DARK_GRAY))
            add_text(slide, v, Inches(x), Inches(1.9 + r * 0.39), Inches(w), Inches(0.38),
                     size=10.5, bold=(r == 0 or "↑" in v or "↓" in v),
                     color=fc, align=PP_ALIGN.CENTER)

    # Interpretação
    add_rect(slide, Inches(5.2), Inches(4.55), Inches(7.9), Inches(2.6),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.8)
    add_text(slide, "O que aprendemos com o V2",
             Inches(5.4), Inches(4.6), Inches(7.5), Inches(0.42),
             size=13, bold=True, color=C_DARK_BLUE)
    add_bullets(slide, [
        "LR menor aproxima o modelo do 'limite entre os dois lados' da fronteira de decisão",
        "FN caem de 18 → 9: o modelo acerta mais pneumonias ambíguas",
        "FP sobem de 38 → 53: o modelo fica mais 'paranoico' com normais",
        "Iterações adicionais de XAI não eliminam o tradeoff — apenas o reposicionam",
        "Conclusão: escolha do modelo depende do custo clínico FP vs FN",
    ], Inches(5.4), Inches(5.1), Inches(7.5), Inches(2.0), size=11)


def slide_tradeoff(prs: Presentation, tradeoff_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Tradeoff FP × FN: Qual Modelo Para Qual Contexto?",
                "A escolha depende do custo clínico de cada tipo de erro")

    slide.shapes.add_picture(tradeoff_buf, Inches(0.3), Inches(1.35), Inches(6.8), Inches(4.8))

    # Painel de guia clínico
    add_rect(slide, Inches(7.4), Inches(1.35), Inches(5.7), Inches(5.7),
             fill=RGBColor(0xF8, 0xF8, 0xF8), line=C_ACCENT, line_width_pt=1.0)
    add_text(slide, "Guia de Escolha Clínica",
             Inches(7.6), Inches(1.4), Inches(5.3), Inches(0.45),
             size=14, bold=True, color=C_ACCENT)

    contexts = [
        ("Triagem de Emergência",
         "XAI-Guided V2",
         "Minimizar FN é crítico — perder uma\npneumonia é mais grave que alarme falso.\nFN=9 é o menor.",
         "#17BECF"),
        ("Triagem Ambulatorial / Rotina",
         "XAI-Guided V1 ★",
         "Melhor equilíbrio. FP=38, FN=18.\nMenor taxa de erro total (56 vs 84).\nAcurácia e F1 superiores.",
         "#2CA02C"),
        ("Rastreamento Populacional",
         "Head Only",
         "Alta precisão (0,9235). Poucos\nalarmes falsos. Mas 52 FN —\naceitável em triagem inicial.",
         "#7B68EE"),
    ]
    for i, (ctx, rec, desc, clr) in enumerate(contexts):
        y = 2.0 + i * 1.7
        clr_rgb = RGBColor(int(clr[1:3], 16), int(clr[3:5], 16), int(clr[5:7], 16))
        add_rect(slide, Inches(7.6), Inches(y), Inches(5.3), Inches(1.55),
                 fill=RGBColor(0xF0, 0xF0, 0xF0), line=clr_rgb, line_width_pt=1.2)
        add_text(slide, f"Contexto: {ctx}", Inches(7.7), Inches(y + 0.03),
                 Inches(5.1), Inches(0.35), size=11, bold=True, color=clr_rgb)
        add_text(slide, f"Modelo recomendado: {rec}",
                 Inches(7.7), Inches(y + 0.38), Inches(5.1), Inches(0.32),
                 size=11, bold=True, color=C_DARK_BLUE)
        add_text(slide, desc, Inches(7.7), Inches(y + 0.72), Inches(5.1), Inches(0.78),
                 size=10, color=C_DARK_GRAY)

    add_text(slide, "FP = alarme falso (normal → pneumonia)  |  FN = pneumonia perdida (pneumonia → normal)",
             Inches(0.3), Inches(6.35), Inches(7.0), Inches(0.45),
             size=10, italic=True, color=C_MID_GRAY, align=PP_ALIGN.CENTER)


def slide_atencao_gradcam(prs: Presentation, overlap_buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Análise de Atenção Grad-CAM++: Como Evoluiu?",
                "overlap_score nos erros — fração da ativação FORA da máscara pulmonar proxy")

    slide.shapes.add_picture(overlap_buf, Inches(0.3), Inches(1.35), Inches(7.0), Inches(4.8))

    # Tabela de erros por modelo
    add_rect(slide, Inches(7.6), Inches(1.35), Inches(5.5), Inches(5.7),
             fill=C_LIGHT_BLUE, line=C_MED_BLUE, line_width_pt=0.8)
    add_text(slide, "Distribuição dos Erros por Modelo",
             Inches(7.8), Inches(1.4), Inches(5.2), Inches(0.42),
             size=13, bold=True, color=C_DARK_BLUE)

    err_rows = [
        ("Modelo", "FP", "FN", "Erros", "partial", "distracted"),
        ("Full Finetune", "80", "4",  "84", "83", "1"),
        ("XAI-Guided V1", "38", "18", "56", "52", "4"),
        ("XAI-Guided V2", "53", "9",  "62", "60", "2"),
    ]
    ecws = [2.0, 0.6, 0.6, 0.85, 0.85, 1.1]
    ecxs = [7.65]
    for w in ecws[:-1]:
        ecxs.append(ecxs[-1] + w)

    bg_rows = [C_DARK_BLUE, RGBColor(0xEE, 0xEE, 0xFF),
               RGBColor(0xD5, 0xFF, 0xD5), RGBColor(0xE0, 0xFF, 0xFF)]

    for r, row in enumerate(err_rows):
        for c, (v, w, x) in enumerate(zip(row, ecws, ecxs)):
            bg = bg_rows[r] if r == 0 else bg_rows[r]
            add_rect(slide, Inches(x), Inches(1.88 + r * 0.45), Inches(w), Inches(0.43),
                     fill=bg, line=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.3)
            fc = C_WHITE if r == 0 else C_DARK_GRAY
            add_text(slide, v, Inches(x), Inches(1.88 + r * 0.45), Inches(w), Inches(0.43),
                     size=10, bold=(r == 0), color=fc, align=PP_ALIGN.CENTER)

    add_bullets(slide, [
        "Todos os erros remanescentes estão na faixa 'partial'",
        "  → ativação na cavidade torácica mas extrapolando a máscara central",
        "",
        "XAI-V1 reduziu erros totais: 84 → 56 (−33%)",
        "XAI-V2: erros sobem levemente (56 → 62)",
        "  mas a distribuição FP/FN muda a favor de menos FN",
        "",
        "Nenhum modelo elimina completamente os casos partial",
        "  → Segmentação pulmonar real permitiria overlap mais preciso",
    ], Inches(7.8), Inches(4.05), Inches(5.2), Inches(2.8), size=10.5)


def slide_licoes(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_header(slide, "Lições Aprendidas e Contribuições do Projeto",
                "O que este trabalho mostra que ainda não estava claro na literatura")

    cards = [
        ("XAI como guia de treino",
         ["Grad-CAM++ não é só diagnóstico pós-treino",
          "É um sinal de qualidade por amostra",
          "Sem anotações manuais adicionais",
          "Diferencial vs literatura (Šefčík 2023)"],
         "#2CA02C"),
        ("Iterações XAI são configuráveis",
         ["V1 → equilibra FP/FN (+4,5% acc)",
          "V2 → reposiciona para menos FN",
          "O tradeoff persiste, mas pode ser ajustado",
          "Calibração de threshold complementa"],
         "#1A6CB0"),
        ("Pipeline end-to-end reproduzível",
         ["Dado bruto → 4 checkpoints → figuras",
          "Todos os scripts versionados",
          "Artigo LaTeX com resultados reais",
          "Sem dependências externas de anotação"],
         "#E86A2B"),
        ("Limitações honestas",
         ["Máscara pulmonar proxy ≠ segmentação real",
          "Split estático (sem k-fold)",
          "Dataset único (Kaggle Chest X-Ray)",
          "CPU only — sem GPU"],
         "#7B68EE"),
    ]

    for i, (title, bullets, color_hex) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.55
        y = 1.38 + row * 3.0
        clr = RGBColor(int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16))
        add_rect(slide, Inches(x), Inches(y), Inches(6.3), Inches(2.8),
                 fill=RGBColor(0xF5, 0xF5, 0xF5), line=clr, line_width_pt=1.5)
        add_rect(slide, Inches(x), Inches(y), Inches(6.3), Inches(0.5), fill=clr)
        add_text(slide, title, Inches(x + 0.1), Inches(y), Inches(6.1), Inches(0.5),
                 size=13, bold=True, color=C_WHITE)
        add_bullets(slide, bullets, Inches(x + 0.1), Inches(y + 0.55),
                    Inches(6.1), Inches(2.1), size=11)


def slide_conclusao_final(prs: Presentation, models_data: list[tuple]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_DARK_BLUE)
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(4.4),
             fill=RGBColor(0x0A, 0x2A, 0x52))

    add_text(slide, "Conclusão Final",
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6),
             size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "XAI-Guided Fine-Tuning: Grad-CAM++ como ferramenta de melhoria de modelo",
             Inches(0.5), Inches(0.78), Inches(12.3), Inches(0.45),
             size=16, italic=True, color=RGBColor(0xC0, 0xD8, 0xF0), align=PP_ALIGN.CENTER)

    # Quadro de resultados
    add_rect(slide, Inches(0.3), Inches(1.32), Inches(7.8), Inches(4.25),
             fill=RGBColor(0x0D, 0x33, 0x60), line=C_MED_BLUE, line_width_pt=0.8)
    add_text(slide, "Os 4 Modelos em Números",
             Inches(0.5), Inches(1.38), Inches(7.5), Inches(0.42),
             size=14, bold=True, color=C_LIGHT_BLUE)

    rows_txt = [
        ("Head Only",      "87,18%", "0,9300", "FP=28  FN=52", "#7B68EE"),
        ("Full Finetune",  "86,54%", "0,9693", "FP=80  FN=4",  "#4A9CE8"),
        ("XAI-Guided V1 ★","91,03%", "0,9663", "FP=38  FN=18", "#2CA02C"),
        ("XAI-Guided V2",  "90,06%", "0,9655", "FP=53  FN=9",  "#17BECF"),
    ]
    for i, (nm, acc, auc_, fp_fn, clr_hex) in enumerate(rows_txt):
        y = 1.88 + i * 0.7
        clr = RGBColor(int(clr_hex[1:3], 16), int(clr_hex[3:5], 16), int(clr_hex[5:7], 16))
        add_text(slide, f"{nm}:", Inches(0.5), Inches(y), Inches(2.8), Inches(0.55),
                 size=12, bold=True, color=clr)
        add_text(slide, f"acc={acc}  AUC={auc_}  {fp_fn}",
                 Inches(3.2), Inches(y), Inches(4.8), Inches(0.55),
                 size=12, color=C_WHITE)

    # Mensagem central
    add_rect(slide, Inches(8.35), Inches(1.32), Inches(4.7), Inches(4.25),
             fill=RGBColor(0x0D, 0x33, 0x60), line=C_GREEN, line_width_pt=1.5)
    add_text(slide, "A mensagem central",
             Inches(8.5), Inches(1.38), Inches(4.4), Inches(0.42),
             size=14, bold=True, color=C_GREEN)
    add_bullets(slide, [
        "XAI não é só auditoria —",
        "é melhoria de modelo",
        "",
        "−52,5% FP com V1",
        "−50% FN com V2",
        "",
        "Sem anotações manuais",
        "Sem nova arquitetura",
        "Sem novos dados",
        "",
        "Ciclo iterável e",
        "configurável clinicamente",
    ], Inches(8.5), Inches(1.85), Inches(4.4), Inches(3.5), size=12,
    color=C_WHITE, bullet="")

    # Trabalhos futuros
    add_rect(slide, Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.5),
             fill=RGBColor(0x0D, 0x33, 0x60), line=C_ACCENT, line_width_pt=0.8)
    add_text(slide, "Trabalhos futuros:",
             Inches(0.5), Inches(5.73), Inches(2.2), Inches(0.4),
             size=12, bold=True, color=C_YELLOW)
    add_text(slide,
             "Segmentação pulmonar real (U-Net)  ·  Calibração de threshold  ·  "
             "Validação cruzada k-fold  ·  Extensão multiclasse  ·  Discussão clínica com especialista",
             Inches(2.6), Inches(5.75), Inches(10.2), Inches(0.4),
             size=11, color=C_WHITE)
    add_text(slide,
             "Obrigado — perguntas?",
             Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.55),
             size=22, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)


# ─── Main ────────────────────────────────────────────────────────────────────

def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path,
                        default=ROOT / "apresentacao_final.pptx")
    args = parser.parse_args()

    # Carrega métricas
    m_head  = load_json(RESULTS_DIR / "metrics_head_only.json")
    m_full  = load_json(RESULTS_DIR / "metrics_full_finetune.json")
    m_v1    = load_json(RESULTS_DIR / "metrics_xai_guided.json")
    m_v2    = load_json(RESULTS_DIR / "metrics_xai_guided_v2.json") \
              if (RESULTS_DIR / "metrics_xai_guided_v2.json").is_file() else None

    # (nome, FP, FN, AUC, F1, acurácia)
    models_data = [
        ("Head Only",     m_head["test"]["confusion_matrix"][0][1],
                          m_head["test"]["confusion_matrix"][1][0],
                          m_head["test"]["auc_roc"], m_head["test"]["f1"],
                          m_head["test"]["accuracy_percent"]),
        ("Full Finetune", m_full["test"]["confusion_matrix"][0][1],
                          m_full["test"]["confusion_matrix"][1][0],
                          m_full["test"]["auc_roc"], m_full["test"]["f1"],
                          m_full["test"]["accuracy_percent"]),
        ("XAI-Guided V1", m_v1["test"]["confusion_matrix"][0][1],
                          m_v1["test"]["confusion_matrix"][1][0],
                          m_v1["test"]["auc_roc"], m_v1["test"]["f1"],
                          m_v1["test"]["accuracy_percent"]),
    ]
    if m_v2:
        models_data.append((
            "XAI-Guided V2", m_v2["test"]["confusion_matrix"][0][1],
                             m_v2["test"]["confusion_matrix"][1][0],
                             m_v2["test"]["auc_roc"], m_v2["test"]["f1"],
                             m_v2["test"]["accuracy_percent"]
        ))

    print("Gerando figuras...")
    pipeline_buf  = make_pipeline_png()
    bar_buf       = make_4models_bar_png(models_data)
    evo_buf       = make_metrics_evolution_png(models_data)
    tradeoff_buf  = make_tradeoff_diagram_png(models_data)
    overlap_buf   = make_overlap_comparison_png()

    print("Montando slides...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)                                          # 1
    slide_ciclo_completo(prs, pipeline_buf)                  # 2
    slide_promessa_vs_entregue(prs)                          # 3
    slide_pipeline_final(prs, pipeline_buf)                  # 4
    slide_resultados_tabela(prs, models_data)                # 5
    slide_evolucao_metricas(prs, evo_buf, bar_buf)           # 6
    slide_xai_metodo(prs)                                    # 7
    slide_v2_deep_dive(prs)                                  # 8
    slide_tradeoff(prs, tradeoff_buf)                        # 9
    slide_atencao_gradcam(prs, overlap_buf)                  # 10
    slide_licoes(prs)                                        # 11
    slide_conclusao_final(prs, models_data)                  # 12

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print(f"\nApresentação final salva em: {args.out}")
    print(f"Total de slides: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
